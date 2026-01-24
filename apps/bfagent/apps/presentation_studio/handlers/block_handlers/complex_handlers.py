"""Complex Block Handlers - Vertical Boxes, Function List"""
from pptx.util import Inches, Pt
from .base_handler import BaseBlockHandler


class VerticalBoxesHandler(BaseBlockHandler):
    """Handler for vertical boxes blocks"""
    
    def render(self, content, styling, layout):
        """Render vertical boxes"""
        title = content.get('title', '')
        boxes = content.get('boxes', [])
        
        y_start = self.current_y
        
        # Title
        if title:
            title_box = self.slide.shapes.add_textbox(
                Inches(5), 
                Inches(y_start), 
                Inches(4.5), 
                Inches(0.5)
            )
            title_box.text_frame.text = title
            if title_box.text_frame.paragraphs:
                p = title_box.text_frame.paragraphs[0]
                p.font.size = self._get_font_size(styling, 'title_size', 22)
                p.font.bold = True
            y_start += 0.7
        
        # Boxes
        top = y_start + 0.5
        for box_data in boxes:
            box = self.slide.shapes.add_textbox(
                Inches(5), 
                Inches(top), 
                Inches(4.5), 
                Inches(1)
            )
            icon = box_data.get('icon', '')
            system = box_data.get('system', '')
            outcome = box_data.get('outcome', '')
            
            text = f"{icon} {system} → {outcome}"
            box.text_frame.text = text
            if box.text_frame.paragraphs:
                p = box.text_frame.paragraphs[0]
                p.font.size = self._get_font_size(styling, 'system_size', 16)
            
            top += 1.2
        
        return top


class FunctionListHandler(BaseBlockHandler):
    """Handler for function list blocks"""
    
    def render(self, content, styling, layout):
        """Render function list"""
        title = content.get('title', '')
        functions = content.get('functions', [])
        
        y_start = self.current_y
        
        # Title
        if title:
            title_box = self.slide.shapes.add_textbox(
                Inches(0.5), 
                Inches(y_start), 
                Inches(9), 
                Inches(0.5)
            )
            title_box.text_frame.text = title
            if title_box.text_frame.paragraphs:
                p = title_box.text_frame.paragraphs[0]
                p.font.size = self._get_font_size(styling, 'title_size', 22)
                p.font.bold = True
            y_start += 0.7
        
        # Functions in 2 columns
        col_width = Inches(4.5)
        row_height = 1.2
        items_per_col = 3
        
        for idx, func in enumerate(functions[:6]):  # Max 6 items
            col = idx // items_per_col
            row = idx % items_per_col
            
            left = Inches(0.5 + col * 5)
            top = Inches(y_start + 0.5 + row * row_height)
            
            box = self.slide.shapes.add_textbox(left, top, col_width, Inches(row_height))
            frame = box.text_frame
            frame.word_wrap = True
            
            icon = func.get('icon', '')
            number = func.get('number', idx + 1)
            name = func.get('name', '')
            desc = func.get('description', '')
            
            text = f"{icon} {number}. {name}: {desc}"
            frame.text = text
            
            if frame.paragraphs:
                p = frame.paragraphs[0]
                p.font.size = Pt(14)
        
        return y_start + 4
