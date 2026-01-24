"""Layout Block Handlers - Two Column, Definition Box"""
from pptx.util import Inches, Pt
from .base_handler import BaseBlockHandler


class DefinitionBoxHandler(BaseBlockHandler):
    """Handler for definition box blocks"""
    
    def render(self, content, styling, layout):
        """Render definition box"""
        term = content.get('term', '')
        definition = content.get('definition', '')
        metaphor = content.get('metaphor', '')
        
        # Calculate safe height
        desired_height = 3.0
        safe_height = self._calculate_safe_height(desired_height)
        
        box = self.slide.shapes.add_textbox(
            Inches(1), 
            Inches(self.current_y), 
            Inches(8), 
            Inches(safe_height)
        )
        box.text_frame.word_wrap = True
        
        # Term
        if term:
            p = box.text_frame.paragraphs[0]
            p.text = term
            p.font.size = self._get_font_size(styling, 'term_size', 28)
            p.font.bold = True
        
        # Definition
        if definition:
            p = box.text_frame.add_paragraph()
            p.text = definition
            p.font.size = self._get_font_size(styling, 'definition_size', 16)
        
        # Metaphor
        if metaphor:
            p = box.text_frame.add_paragraph()
            p.text = f"\n{metaphor}"
            p.font.size = self._get_font_size(styling, 'metaphor_size', 18)
            p.font.italic = True
        
        return self.current_y + safe_height + 0.3


class TwoColumnComparisonHandler(BaseBlockHandler):
    """Handler for two column comparison blocks"""
    
    def render(self, content, styling, layout):
        """Render two column comparison"""
        title = content.get('title', '')
        col_left = content.get('column_left', {})
        col_right = content.get('column_right', {})
        
        y_start = self.current_y
        
        # Title
        if title:
            title_box = self.slide.shapes.add_textbox(
                Inches(0.5), 
                Inches(y_start), 
                Inches(9), 
                Inches(0.5)
            )
            title_box.text_frame.text = title
            if title_box.text_frame.paragraphs:
                p = title_box.text_frame.paragraphs[0]
                p.font.size = self._get_font_size(styling, 'title_size', 22)
                p.font.bold = True
            y_start += 0.7
        
        # Calculate safe height for columns
        desired_col_height = 4.5
        safe_col_height = self._calculate_safe_height(desired_col_height + (y_start - self.current_y))
        
        # Left column
        left_box = self.slide.shapes.add_textbox(
            Inches(0.5), 
            Inches(y_start), 
            Inches(4.5), 
            Inches(safe_col_height)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        # Left heading
        if 'heading' in col_left:
            p = left_frame.paragraphs[0]
            p.text = col_left['heading']
            p.font.size = self._get_font_size(styling, 'heading_size', 18)
            p.font.bold = True
        
        # Left items (simplified)
        items_left = col_left.get('items', [])
        for item in items_left[:3]:  # Limit to 3 items
            p = left_frame.add_paragraph()
            term = item.get('term', '')
            desc = item.get('description', '')
            p.text = f"• {term}: {desc}"
            p.font.size = Pt(14)
            p.level = 1
        
        # Right column
        right_box = self.slide.shapes.add_textbox(
            Inches(5.5), 
            Inches(y_start), 
            Inches(4.5), 
            Inches(safe_col_height)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        # Right heading
        if 'heading' in col_right:
            p = right_frame.paragraphs[0]
            p.text = col_right['heading']
            p.font.size = self._get_font_size(styling, 'heading_size', 18)
            p.font.bold = True
        
        # Right items (simplified)
        items_right = col_right.get('items', [])
        for item in items_right[:3]:  # Limit to 3 items
            p = right_frame.add_paragraph()
            term = item.get('term', '')
            desc = item.get('description', '')
            p.text = f"• {term}: {desc}"
            p.font.size = Pt(14)
            p.level = 1
        
        return y_start + safe_col_height + 0.3
