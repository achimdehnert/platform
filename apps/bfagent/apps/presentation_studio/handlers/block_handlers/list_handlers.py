"""List Block Handlers - Bullet Lists, Objectives"""
from pptx.util import Inches, Pt
from .base_handler import BaseBlockHandler


class BulletListHandler(BaseBlockHandler):
    """Handler for bullet list blocks"""
    
    def render(self, content, styling, layout):
        """Render bullet list"""
        title = content.get('title', '')
        items = content.get('items', [])
        
        width = Inches(layout.get('width_percent', 80) / 10)
        left = Inches((10 - width.inches) / 2)
        
        box = self.slide.shapes.add_textbox(
            left, 
            Inches(self.current_y), 
            width, 
            Inches(4)
        )
        
        # Title
        if title:
            p = box.text_frame.paragraphs[0]
            p.text = title
            p.font.size = self._get_font_size(styling, 'title_size', 22)
            p.font.bold = True
        
        # Items
        for item in items:
            p = box.text_frame.add_paragraph()
            p.text = item
            p.font.size = self._get_font_size(styling, 'item_size', 16)
            p.level = 1
        
        return self.current_y + 4


class ObjectivesBoxHandler(BaseBlockHandler):
    """Handler for objectives box blocks"""
    
    def render(self, content, styling, layout):
        """Render objectives box"""
        heading = content.get('heading', '')
        objectives = content.get('objectives', [])
        
        box = self.slide.shapes.add_textbox(
            Inches(0.5), 
            Inches(self.current_y), 
            Inches(9), 
            Inches(2)
        )
        
        # Heading
        if heading:
            p = box.text_frame.paragraphs[0]
            p.text = heading
            p.font.size = self._get_font_size(styling, 'heading_size', 18)
            p.font.bold = True
        
        # Objectives
        for obj in objectives:
            p = box.text_frame.add_paragraph()
            checkmark = styling.get('checkmark', '✓')
            p.text = f"{checkmark} {obj}"
            p.font.size = self._get_font_size(styling, 'objective_size', 14)
            p.level = 1
        
        return self.current_y + 2
