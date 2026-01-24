"""Text Block Handlers - Title, Quote, TextBox"""
from pptx.util import Inches, Pt
from .base_handler import BaseBlockHandler


class TitleHandler(BaseBlockHandler):
    """Handler for title blocks"""
    
    def render(self, content, styling, layout):
        """Render title block"""
        box = self.slide.shapes.add_textbox(
            Inches(0.5), 
            Inches(self.current_y), 
            Inches(9), 
            Inches(1.5)
        )
        box.text_frame.text = content.get('main_title', '')
        
        if box.text_frame.paragraphs:
            p = box.text_frame.paragraphs[0]
            p.font.size = self._get_font_size(styling, 'title_size', 44)
            p.font.bold = True
        
        return self.current_y + 1.5


class QuoteHandler(BaseBlockHandler):
    """Handler for quote blocks"""
    
    def render(self, content, styling, layout):
        """Render quote block"""
        box = self.slide.shapes.add_textbox(
            Inches(2), 
            Inches(self.current_y), 
            Inches(6), 
            Inches(2)
        )
        box.text_frame.text = f'"{content.get("quote", "")}"'
        
        if box.text_frame.paragraphs:
            p = box.text_frame.paragraphs[0]
            p.font.size = self._get_font_size(styling, 'quote_size', 22)
            p.font.italic = True
        
        return self.current_y + 2


class TextBoxHandler(BaseBlockHandler):
    """Handler for text box blocks"""
    
    def render(self, content, styling, layout):
        """Render text box"""
        width = Inches(4.5)
        position = layout.get('position', '')
        left = Inches(0.5 if 'left' in position else 5)
        
        box = self.slide.shapes.add_textbox(
            left, 
            Inches(self.current_y), 
            width, 
            Inches(4)
        )
        
        # Heading
        if 'heading' in content:
            p = box.text_frame.paragraphs[0]
            p.text = content['heading']
            p.font.size = self._get_font_size(styling, 'heading_size', 24)
            p.font.bold = True
        
        # Body
        if 'body' in content:
            p = box.text_frame.add_paragraph()
            p.text = content['body']
            p.font.size = self._get_font_size(styling, 'body_size', 16)
        
        return self.current_y + 4
