"""
Presentation Enhancement Handler
Handles PPTX content enhancement with concept-based slides
"""

import os
import logging
from typing import Dict, List, Optional
from pathlib import Path
from datetime import datetime

from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from django.core.files.base import ContentFile
from django.conf import settings

logger = logging.getLogger(__name__)


class EnhanceContentHandler:
    """
    Handler for enhancing PowerPoint presentations with concept-based slides
    """

    # Concept templates by type
    CONCEPT_TEMPLATES = {
        'medical': [
            {'title': 'Patient Case Overview', 'type': 'title_content'},
            {'title': 'Diagnosis & Assessment', 'type': 'two_column'},
            {'title': 'Treatment Plan', 'type': 'bullet_list'},
        ],
        'business': [
            {'title': 'Executive Summary', 'type': 'title_content'},
            {'title': 'Market Analysis', 'type': 'two_column'},
            {'title': 'Strategic Recommendations', 'type': 'bullet_list'},
        ],
        'scientific': [
            {'title': 'Research Objectives', 'type': 'title_content'},
            {'title': 'Methodology', 'type': 'two_column'},
            {'title': 'Results & Conclusions', 'type': 'bullet_list'},
        ],
        'technical': [
            {'title': 'Technical Overview', 'type': 'title_content'},
            {'title': 'Architecture & Design', 'type': 'two_column'},
            {'title': 'Implementation Details', 'type': 'bullet_list'},
        ],
    }

    def extract_slide_count(self, pptx_path: str) -> int:
        """
        Extract the number of slides from a PPTX file
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Number of slides
        """
        try:
            prs = PptxPresentation(pptx_path)
            return len(prs.slides)
        except Exception as e:
            logger.error(f"Error extracting slide count: {str(e)}")
            return 0

    def extract_table_of_contents(self, pptx_path: str) -> List[Dict]:
        """
        Extract table of contents from a PPTX file
        Returns slide numbers, titles, and text previews
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            List of dicts with slide_number, title, preview_text
        """
        try:
            prs = PptxPresentation(pptx_path)
            toc = []
            
            for idx, slide in enumerate(prs.slides, start=1):
                slide_info = {
                    'slide_number': idx,
                    'title': self._extract_slide_title(slide, idx),
                    'preview_text': self._extract_slide_preview(slide),
                    'shape_count': len(slide.shapes),
                }
                toc.append(slide_info)
            
            return toc
            
        except Exception as e:
            logger.error(f"Error extracting table of contents: {str(e)}")
            return []

    def _extract_slide_title(self, slide, slide_number: int = None) -> str:
        """Extract title from a slide"""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
            # Use sequential number instead of slide_id
            if slide_number:
                return f"Slide {slide_number}"
            return "Untitled Slide"
        except Exception:
            return "Untitled Slide"

    def _extract_slide_preview(self, slide, max_length: int = 150) -> str:
        """Extract text preview from a slide (first few lines)"""
        try:
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())
            
            full_text = " ".join(text_parts)
            if len(full_text) > max_length:
                return full_text[:max_length] + "..."
            return full_text if full_text else "(No text content)"
            
        except Exception:
            return "(Preview unavailable)"

    def enhance_presentation(
        self,
        original_path: str,
        enhancement_type: str,
        concepts: List[Dict],
        mode: str = 'append'
    ) -> Dict:
        """
        Enhance a presentation with concept-based slides
        
        Args:
            original_path: Path to original PPTX file
            enhancement_type: Type of enhancement (medical, business, etc.)
            concepts: List of concepts to add
            mode: Enhancement mode ('append' or 'smart')
            
        Returns:
            Dict with success status and enhanced file path
        """
        try:
            # Load original presentation
            prs = PptxPresentation(original_path)
            slides_before = len(prs.slides)
            
            # Get concept templates
            templates = self.CONCEPT_TEMPLATES.get(
                enhancement_type,
                self.CONCEPT_TEMPLATES['technical']
            )
            
            # Add concept slides
            slides_added = 0
            for i, concept in enumerate(concepts):
                template = templates[i % len(templates)]
                self._add_concept_slide(prs, concept, template)
                slides_added += 1
            
            # Save enhanced presentation
            enhanced_path = self._get_enhanced_path(original_path)
            prs.save(enhanced_path)
            
            slides_after = len(prs.slides)
            
            return {
                'success': True,
                'enhanced_path': enhanced_path,
                'slides_before': slides_before,
                'slides_after': slides_after,
                'slides_added': slides_added,
            }
            
        except Exception as e:
            logger.error(f"Error enhancing presentation: {str(e)}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'slides_before': 0,
                'slides_after': 0,
                'slides_added': 0,
            }

    def _find_best_layout(self, prs: PptxPresentation, layout_type: str):
        """
        Find the best matching layout from the existing presentation
        
        Args:
            prs: Presentation object
            layout_type: Type of layout needed ('title_content', 'two_column', 'bullet_list')
        
        Returns:
            Best matching slide layout
        """
        # Try to find a layout with both title and content placeholders
        for layout in prs.slide_layouts:
            try:
                has_title = any(shape.placeholder_format.type == 1 for shape in layout.placeholders)
                has_body = any(shape.placeholder_format.type == 2 for shape in layout.placeholders)
                
                if has_title and has_body:
                    return layout
            except:
                continue
        
        # Fallback: Use first available layout
        return prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[1]
    
    def _add_concept_slide(
        self,
        prs: PptxPresentation,
        concept: Dict,
        template: Dict
    ):
        """
        Add a single concept slide to the presentation
        
        Args:
            prs: Presentation object
            concept: Concept data
            template: Template configuration
        """
        # Check if this is a markdown-parsed slide with full structure
        if 'slide_data' in concept:
            self._add_markdown_slide(prs, concept)
            return
        
        # Get template type (support both template['type'] and concept['slide_type'])
        template_type = concept.get('slide_type') or template.get('type', 'title_content')
        
        if template_type == 'title_content':
            self._add_title_content_slide(prs, concept, template)
        elif template_type == 'two_column':
            self._add_two_column_slide(prs, concept, template)
        elif template_type == 'bullet_list':
            self._add_bullet_list_slide(prs, concept, template)
        else:
            self._add_title_content_slide(prs, concept, template)
    
    def _add_markdown_slide(self, prs: PptxPresentation, concept: Dict):
        """
        Add a slide from markdown-parsed data using MarkdownToSlidesHandler
        
        Args:
            prs: Presentation object
            concept: Concept with slide_data
        """
        try:
            from .markdown_to_slides_handler import MarkdownToSlidesHandler
            from .markdown_slide_parser import SlideContent
            
            slide_data = concept.get('slide_data', {})
            
            # Create SlideContent object
            slide_content = SlideContent(
                slide_number=slide_data.get('slide_number', 1),
                title=concept.get('title', ''),
                navigation=slide_data.get('navigation'),
                headline=concept.get('content', ''),
                quote=slide_data.get('quote'),
                quote_author=slide_data.get('quote_author'),
                content_blocks=slide_data.get('content_blocks', []),
                sources=slide_data.get('sources', []),
                visual_notes=slide_data.get('visual_notes', [])
            )
            
            # Use MarkdownToSlidesHandler to create the slide
            handler = MarkdownToSlidesHandler(prs)
            handler._create_slide(slide_content)
            
            logger.info(f"Created markdown slide: {concept.get('title')}")
            
        except Exception as e:
            logger.error(f"Error creating markdown slide: {e}")
            # Fallback to basic slide
            self._add_title_content_slide(prs, concept, {})

    def _add_title_content_slide(
        self,
        prs: PptxPresentation,
        concept: Dict,
        template: Dict
    ):
        """Add a title and content slide using existing presentation's layout"""
        # Find the best matching layout from the existing presentation
        slide_layout = self._find_best_layout(prs, 'title_content')
        slide = prs.slides.add_slide(slide_layout)
        
        title = concept.get('title', template.get('title', 'Concept'))
        content = concept.get('content', concept.get('description', ''))
        
        # DEBUG LOGGING
        logger.info(f"[TITLE_CONTENT_SLIDE] Title: {title[:50] if title else 'EMPTY'}")
        logger.info(f"[TITLE_CONTENT_SLIDE] Content length: {len(content) if content else 0}")
        logger.info(f"[TITLE_CONTENT_SLIDE] Concept keys: {list(concept.keys())}")
        
        # Try to set title
        if slide.shapes.title:
            slide.shapes.title.text = title
            logger.info(f"[TITLE_CONTENT_SLIDE] Title set successfully")
        else:
            logger.warning(f"[TITLE_CONTENT_SLIDE] No title shape found!")
        
        # ROBUST CONTENT PLACEMENT with multiple fallback strategies
        placed = self._place_text_content(slide, content)
        logger.info(f"[TITLE_CONTENT_SLIDE] Content placement: {placed}")

    def _add_two_column_slide(
        self,
        prs: PptxPresentation,
        concept: Dict,
        template: Dict
    ):
        """Add a two-column slide using existing presentation's layout"""
        # Find the best matching layout
        slide_layout = self._find_best_layout(prs, 'two_column')
        slide = prs.slides.add_slide(slide_layout)
        
        title = concept.get('title', template.get('title', 'Concept'))
        left_content = concept.get('left', concept.get('content', ''))
        right_content = concept.get('right', '')
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Combine content for single placeholder
        combined_content = left_content
        if right_content:
            combined_content += "\n\n" + right_content
        
        # ROBUST CONTENT PLACEMENT
        self._place_text_content(slide, combined_content)

    def _add_bullet_list_slide(
        self,
        prs: PptxPresentation,
        concept: Dict,
        template: Dict
    ):
        """Add a bullet list slide using existing presentation's layout"""
        # Find the best matching layout
        slide_layout = self._find_best_layout(prs, 'bullet_list')
        slide = prs.slides.add_slide(slide_layout)
        
        title = concept.get('title', template.get('title', 'Concept'))
        content = concept.get('content', '')
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # ROBUST CONTENT PLACEMENT
        self._place_text_content(slide, content)

    def _place_text_content(self, slide, content: str) -> bool:
        """
        Robustly place text content into slide with multiple fallback strategies
        
        Args:
            slide: Slide object
            content: Text content to place
            
        Returns:
            True if placement succeeded, False otherwise
        """
        if not content:
            return True  # Nothing to place
        
        # STRATEGY 1: Try body placeholder (type 2)
        try:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # BODY
                    shape.text = content
                    logger.info("Content placed in body placeholder")
                    return True
        except Exception as e:
            logger.debug(f"Body placeholder strategy failed: {e}")
        
        # STRATEGY 2: Try any text frame placeholder (excluding title)
        try:
            for shape in slide.placeholders:
                # Skip title placeholder
                if shape.placeholder_format.type == 1:
                    continue
                if hasattr(shape, 'text_frame'):
                    shape.text = content
                    logger.info(f"Content placed in placeholder type {shape.placeholder_format.type}")
                    return True
        except Exception as e:
            logger.debug(f"Text frame placeholder strategy failed: {e}")
        
        # STRATEGY 3: Try any shape with text_frame (non-placeholder)
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                    shape.text = content
                    logger.info("Content placed in text frame shape")
                    return True
        except Exception as e:
            logger.debug(f"Text frame shape strategy failed: {e}")
        
        # STRATEGY 4: Add new text box as last resort
        try:
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = content
            logger.warning("Content placed in new textbox (fallback strategy)")
            return True
        except Exception as e:
            logger.error(f"All content placement strategies failed: {e}")
            return False
    
    def _get_enhanced_path(self, original_path: str) -> str:
        """
        Generate path for enhanced file
        
        Args:
            original_path: Path to original file
            
        Returns:
            Path for enhanced file
        """
        path = Path(original_path)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        enhanced_name = f"{path.stem}_enhanced_{timestamp}{path.suffix}"
        return str(path.parent / enhanced_name)

    def get_available_concepts(self, enhancement_type: str) -> List[Dict]:
        """
        Get available concept templates for a given enhancement type
        
        Args:
            enhancement_type: Type of enhancement
            
        Returns:
            List of concept templates
        """
        return self.CONCEPT_TEMPLATES.get(
            enhancement_type,
            self.CONCEPT_TEMPLATES['technical']
        )
