"""
Markdown to Slides Handler for PPTX Studio
Converts structured markdown content into PowerPoint slides
"""

import logging
from typing import Dict, List, Optional
from pptx import Presentation as PptxPresentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from .markdown_slide_parser import MarkdownSlideParser, SlideContent
from .slide_layout_handler import SlideLayoutHandler

logger = logging.getLogger(__name__)


class MarkdownToSlidesHandler:
    """
    Converts parsed markdown slides to PowerPoint presentation
    
    Features:
    - Intelligent layout selection based on content
    - Text formatting (bold, italic, bullets)
    - Code block handling (formatted as boxes)
    - Quote styling
    - Source citations
    """
    
    def __init__(self, prs, template_collection=None):
        self.prs = prs
        self.template_collection = template_collection
        self.layout_handler = SlideLayoutHandler()
        self.color_scheme = {
            'primary': RGBColor(102, 126, 234),      # Blue
            'secondary': RGBColor(118, 75, 162),     # Purple
            'text': RGBColor(51, 51, 51),           # Dark gray
            'quote': RGBColor(100, 100, 100),       # Medium gray
            'code_bg': RGBColor(245, 245, 245),     # Light gray
        }
    
    def create_slides_from_markdown(
        self,
        markdown_parser: MarkdownSlideParser,
        start_slide: Optional[int] = None,
        end_slide: Optional[int] = None
    ) -> List[int]:
        """
        Create PowerPoint slides from parsed markdown
        
        Args:
            markdown_parser: Parsed markdown content
            start_slide: Optional start slide number (1-indexed)
            end_slide: Optional end slide number (1-indexed)
            
        Returns:
            List of created slide indices in presentation
        """
        created_slides = []
        slides_to_process = markdown_parser.slides
        
        # Filter slides if range specified
        if start_slide or end_slide:
            start = start_slide or 1
            end = end_slide or len(slides_to_process)
            slides_to_process = [
                s for s in slides_to_process 
                if start <= s.slide_number <= end
            ]
        
        for slide_content in slides_to_process:
            try:
                slide_idx = self._create_slide(slide_content)
                created_slides.append(slide_idx)
                logger.info(f"Created slide {slide_content.slide_number}: {slide_content.title}")
            except Exception as e:
                logger.error(f"Error creating slide {slide_content.slide_number}: {e}")
                continue
        
        return created_slides
    
    def _create_slide(self, content: SlideContent) -> int:
        """Create a single PowerPoint slide from content"""
        import logging
        logger = logging.getLogger(__name__)
        
        # STRATEGY: Clone an existing slide with good layout, replace content
        template_slide = self._find_good_template_slide()
        
        if template_slide:
            logger.info(f"Cloning existing slide as template")
            slide = self._clone_slide(template_slide)
            logger.info(f"Cloned slide created with {len(slide.shapes)} shapes")
        else:
            # Fallback: Use layout method
            layout = self._select_layout(content)
            logger.info(f"No template found, using layout: {layout.name if hasattr(layout, 'name') else 'Unknown'}")
            slide = self.prs.slides.add_slide(layout)
            logger.info(f"Slide created with {len(slide.shapes)} shapes")
        
        # Clear and replace title
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = content.headline or content.title
            self._format_title(title)
            logger.info(f"Title replaced: {title.text}")
        else:
            logger.warning("No title shape found!")
        
        # Clear old content and add new
        self._clear_content_shapes(slide)
        
        # Add content based on type
        if content.quote:
            logger.info("Adding quote slide")
            self._add_quote_slide(slide, content)
        elif len(content.content_blocks) > 0:
            logger.info(f"Adding content slide with {len(content.content_blocks)} blocks")
            self._add_content_slide(slide, content)
        else:
            logger.warning("No content to add!")
        
        return len(self.prs.slides) - 1
    
    def _find_good_template_slide(self, template_type='content_slide'):
        """
        Find a good existing slide to use as template
        
        Args:
            template_type: Type of template to find (content_slide, bullet_slide, etc.)
        
        Returns:
            Slide object or None if not found
        """
        import logging
        logger = logging.getLogger(__name__)
        
        # STRATEGY 1: Use TemplateCollection if available
        if self.template_collection:
            return self._get_template_from_collection(template_type)
        
        # STRATEGY 2: Analyze existing slides in current presentation
        return self._find_template_by_analysis()
    
    def _get_template_from_collection(self, template_type='content_slide'):
        """Get template from TemplateCollection"""
        import logging
        from pptx import Presentation as PptxPresentation
        logger = logging.getLogger(__name__)
        
        try:
            templates = self.template_collection.templates
            
            # Try exact match first
            if template_type in templates:
                template_info = templates[template_type]
                slide_idx = template_info['slide_index']
                
                # Load master PPTX and get slide
                master_prs = PptxPresentation(self.template_collection.master_pptx.path)
                
                if slide_idx < len(master_prs.slides):
                    logger.info(f"Using {template_type} from TemplateCollection: {self.template_collection.name}")
                    return master_prs.slides[slide_idx]
            
            # Fallback: Try similar template types
            fallback_types = {
                'content_slide': ['bullet_slide', 'two_column'],
                'bullet_slide': ['content_slide'],
                'quote_slide': ['title_slide'],
            }
            
            for fallback_type in fallback_types.get(template_type, []):
                if fallback_type in templates:
                    template_info = templates[fallback_type]
                    slide_idx = template_info['slide_index']
                    master_prs = PptxPresentation(self.template_collection.master_pptx.path)
                    
                    if slide_idx < len(master_prs.slides):
                        logger.info(f"Using fallback {fallback_type} from TemplateCollection")
                        return master_prs.slides[slide_idx]
            
            logger.warning(f"Template type {template_type} not found in collection")
            return None
            
        except Exception as e:
            logger.error(f"Error loading template from collection: {e}")
            return None
    
    def _find_template_by_analysis(self):
        """Find template by analyzing existing slides (fallback method)"""
        import logging
        logger = logging.getLogger(__name__)
        
        logger.info("Analyzing existing slides for template...")
        
        # Look for slides with title + content (not title-only, not complex)
        for slide in self.prs.slides:
            shape_count = len(slide.shapes)
            has_title = slide.shapes.title is not None
            has_content = False
            
            # Check if has content shapes (textboxes)
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    has_content = True
                    break
            
            # Good template: Has title, has content, not too complex
            if has_title and has_content and 2 <= shape_count <= 5:
                logger.info(f"Found good template slide by analysis with {shape_count} shapes")
                return slide
        
        logger.warning("No good template slide found by analysis")
        return None
    
    def _clone_slide(self, source_slide):
        """Clone an existing slide to create a new one with same layout"""
        # Get the layout of the source slide
        slide_layout = source_slide.slide_layout
        
        # Create new slide with same layout
        new_slide = self.prs.slides.add_slide(slide_layout)
        
        return new_slide
    
    def _clear_content_shapes(self, slide):
        """Clear all content shapes except title"""
        import logging
        logger = logging.getLogger(__name__)
        
        shapes_to_remove = []
        for shape in slide.shapes:
            # Keep title, remove all other text shapes
            if shape.has_text_frame and shape != slide.shapes.title:
                shapes_to_remove.append(shape)
        
        # Remove shapes (must be done separately to avoid iteration issues)
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
        
        logger.info(f"Cleared {len(shapes_to_remove)} content shapes")
    
    def _select_layout(self, content: SlideContent) -> any:
        """Select appropriate layout based on content"""
        import logging
        logger = logging.getLogger(__name__)
        
        # If has quote, use centered layout
        if content.quote and len(content.content_blocks) < 5:
            layout = self._get_layout_by_name(['Title Only', 'Blank'])
            logger.info(f"Quote layout selected: {layout.name if hasattr(layout, 'name') else 'Unknown'}")
            return layout
        
        # If lots of content, use title + content
        if len(content.content_blocks) > 0:
            layout = self._get_layout_by_name(['Title and Content', 'Title Slide', 'Title and Text'])
            logger.info(f"Content layout selected: {layout.name if hasattr(layout, 'name') else 'Unknown'}")
            return layout
        
        # Default: Use first safe layout
        layout = self._get_safe_default_layout()
        logger.info(f"Default layout selected: {layout.name if hasattr(layout, 'name') else 'Unknown'}")
        return layout
    
    def _get_layout_by_name(self, names: List[str]) -> any:
        """Get layout by name (tries multiple names), excluding bad layouts"""
        import logging
        logger = logging.getLogger(__name__)
        
        # Log all available layouts (first time only)
        if not hasattr(self, '_layouts_logged'):
            logger.info("Available layouts:")
            for idx, layout in enumerate(self.prs.slide_layouts):
                logger.info(f"  [{idx}] {layout.name}")
            self._layouts_logged = True
        
        # Try to find matching layout (skip "not 2 use" and similar)
        for layout in self.prs.slide_layouts:
            if layout.name in names:
                # Skip layouts with "not" in name
                if 'not' in layout.name.lower() or 'unused' in layout.name.lower():
                    logger.warning(f"Skipping layout '{layout.name}' (marked as not for use)")
                    continue
                logger.info(f"Found matching layout: {layout.name}")
                return layout
        
        # No match found, use safe default
        logger.warning(f"No layout found matching {names}, using safe default")
        return self._get_safe_default_layout()
    
    def _get_safe_default_layout(self) -> any:
        """Get a safe default layout (avoid 'not 2 use' and bad layouts)"""
        import logging
        logger = logging.getLogger(__name__)
        
        # First try: Find "Blank" layout - we'll add our own text boxes
        for layout in self.prs.slide_layouts:
            name_lower = layout.name.lower()
            if name_lower == 'blank' or name_lower == 'leer':
                logger.info(f"Using Blank layout (will add custom textboxes): {layout.name}")
                return layout
        
        # Second try: Find simple title+content layout
        for layout in self.prs.slide_layouts:
            name_lower = layout.name.lower()
            # Skip known bad layouts
            if 'not' in name_lower or 'unused' in name_lower or 'komplex' in name_lower:
                continue
            # Look for simple layouts
            if any(keyword in name_lower for keyword in ['title and content', 'titel und inhalt', 'content', 'inhalt']):
                logger.info(f"Safe default found: {layout.name}")
                return layout
        
        # Fallback: Use layout 0 (usually title slide)
        logger.warning("No ideal layout found, using layout[0]")
        return self.prs.slide_layouts[0]
    
    def _format_title(self, title_shape):
        """Format title text"""
        if not title_shape.has_text_frame:
            return
        
        text_frame = title_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(40)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.color_scheme['primary']
            paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_quote_slide(self, slide, content: SlideContent):
        """Add a quote-focused slide"""
        # Create text box for quote
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(3)
        
        quote_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = quote_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add quote text
        p = text_frame.paragraphs[0]
        p.text = f'"{content.quote}"'
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = self.color_scheme['quote']
        p.alignment = PP_ALIGN.CENTER
        
        # Add author if present
        if content.quote_author:
            p_author = text_frame.add_paragraph()
            p_author.text = f"— {content.quote_author}"
            p_author.font.size = Pt(20)
            p_author.font.italic = True
            p_author.font.color.rgb = self.color_scheme['text']
            p_author.alignment = PP_ALIGN.RIGHT
            p_author.space_before = Pt(20)
    
    def _add_content_slide(self, slide, content: SlideContent):
        """Add a content-focused slide"""
        # Find content placeholder or create text box
        content_placeholder = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                content_placeholder = shape
                break
        
        if not content_placeholder:
            # Create text box
            left = Inches(0.8)
            top = Inches(1.8)
            width = Inches(8.4)
            height = Inches(4.5)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
        
        text_frame = content_placeholder.text_frame
        text_frame.word_wrap = True
        text_frame.clear()
        
        # Add content blocks
        for idx, block in enumerate(content.content_blocks):
            self._add_content_block(text_frame, block, is_first=(idx == 0))
    
    def _add_content_block(self, text_frame, block: Dict, is_first: bool = False):
        """Add a single content block to text frame"""
        block_type = block.get('type', 'paragraph')
        block_content = block.get('content', '')
        
        if not block_content.strip():
            return
        
        # Create paragraph
        if is_first and len(text_frame.paragraphs) == 1:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        if block_type == 'subheading':
            p.text = block_content
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.color_scheme['secondary']
            p.space_before = Pt(16)
            p.space_after = Pt(8)
        
        elif block_type == 'label':
            p.text = block_content
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.color_scheme['text']
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        
        elif block_type == 'bullet':
            p.text = block_content
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = self.color_scheme['text']
            p.space_after = Pt(4)
        
        elif block_type == 'quote':
            p.text = f'"{block_content}"'
            p.font.size = Pt(18)
            p.font.italic = True
            p.font.color.rgb = self.color_scheme['quote']
            p.space_before = Pt(12)
            p.space_after = Pt(12)
        
        elif block_type == 'code_block':
            # Code blocks as formatted text
            p.text = block_content[:500]  # Limit length
            p.font.name = 'Courier New'
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.space_before = Pt(8)
            p.space_after = Pt(8)
        
        else:  # paragraph
            p.text = block_content
            p.font.size = Pt(16)
            p.font.color.rgb = self.color_scheme['text']
            p.space_after = Pt(8)
    
    def create_summary_slide(
        self,
        markdown_parser: MarkdownSlideParser,
        title: str = "Kapitel Übersicht"
    ) -> int:
        """Create a summary slide with all slide titles"""
        layout = self._get_layout_by_name(['Title and Content', 'Title Slide'])
        slide = self.prs.slides.add_slide(layout)
        
        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._format_title(slide.shapes.title)
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for slide_content in markdown_parser.slides:
            p = text_frame.add_paragraph()
            p.text = f"{slide_content.slide_number}. {slide_content.title}"
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.color_scheme['text']
            p.space_after = Pt(8)
        
        return len(self.prs.slides) - 1


def create_presentation_from_markdown(
    markdown_file_path: str,
    template_path: Optional[str] = None
) -> PptxPresentation:
    """
    Create a complete presentation from markdown file
    
    Args:
        markdown_file_path: Path to markdown file
        template_path: Optional path to PPTX template
        
    Returns:
        PowerPoint presentation object
    """
    from .markdown_slide_parser import parse_markdown_file
    
    # Parse markdown
    parser = parse_markdown_file(markdown_file_path)
    logger.info(f"Parsed {len(parser.slides)} slides from {markdown_file_path}")
    
    # Create presentation
    if template_path:
        prs = PptxPresentation(template_path)
    else:
        prs = PptxPresentation()
    
    # Create handler and generate slides
    handler = MarkdownToSlidesHandler(prs)
    created_slides = handler.create_slides_from_markdown(parser)
    
    logger.info(f"Created {len(created_slides)} slides in presentation")
    
    return prs
