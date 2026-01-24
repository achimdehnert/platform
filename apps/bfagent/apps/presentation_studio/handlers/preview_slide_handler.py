"""
Preview Slide Handler
Manages preview slides before PPTX conversion
"""

import logging
from typing import List, Dict, Optional
from django.utils import timezone
from pptx import Presentation as PptxPresentation

logger = logging.getLogger(__name__)


class PreviewSlideHandler:
    """
    Handler for managing preview slides and their conversion to PPTX
    
    Features:
    - Create preview slides from various sources (Markdown, JSON, PDF)
    - Individual slide conversion
    - Batch conversion
    - Edit preview content
    - Reorder previews
    """
    
    def create_previews_from_markdown(
        self,
        presentation,
        markdown_file_path: str,
        file_name: str = ""
    ) -> List:
        """
        Create preview slides from markdown file
        
        Args:
            presentation: Presentation instance
            markdown_file_path: Path to markdown file
            file_name: Original file name
            
        Returns:
            List of created PreviewSlide objects
        """
        from apps.presentation_studio.models import PreviewSlide
        from apps.presentation_studio.handlers.markdown_slide_parser import parse_markdown_file
        
        try:
            # Parse markdown
            parser = parse_markdown_file(markdown_file_path)
            logger.info(f"Parsed {len(parser.slides)} slides from markdown")
            
            # Get current slide count (for ordering)
            current_count = presentation.slide_count_enhanced or presentation.slide_count_original or 0
            
            # Create preview slides
            preview_slides = []
            for idx, slide_content in enumerate(parser.slides, start=1):
                preview = PreviewSlide.objects.create(
                    presentation=presentation,
                    preview_order=current_count + idx,
                    title=slide_content.title,
                    content_data={
                        'slide_number': slide_content.slide_number,
                        'title': slide_content.title,
                        'headline': slide_content.headline,
                        'navigation': slide_content.navigation,
                        'quote': slide_content.quote,
                        'quote_author': slide_content.quote_author,
                        'content_blocks': slide_content.content_blocks,
                        'sources': slide_content.sources,
                        'visual_notes': slide_content.visual_notes,
                    },
                    source_type='markdown',
                    source_file_name=file_name,
                    status='preview'
                )
                preview_slides.append(preview)
                logger.info(f"Created preview slide {idx}: {slide_content.title}")
            
            return preview_slides
        
        except Exception as e:
            logger.error(f"Error creating previews from markdown: {e}", exc_info=True)
            raise
    
    def convert_preview_to_pptx(self, preview_slide_id: str) -> Dict:
        """
        Convert single preview slide to PPTX
        
        Args:
            preview_slide_id: UUID of preview slide
            
        Returns:
            Result dict with success status and slide number
        """
        from apps.presentation_studio.models import PreviewSlide
        from apps.presentation_studio.handlers.markdown_to_slides_handler import MarkdownToSlidesHandler
        from apps.presentation_studio.handlers.markdown_slide_parser import SlideContent
        
        try:
            preview = PreviewSlide.objects.get(id=preview_slide_id)
            
            if preview.status == 'converted':
                return {
                    'success': False,
                    'error': 'Slide already converted'
                }
            
            presentation = preview.presentation
            
            # Load existing PPTX or create new
            if presentation.enhanced_file and presentation.enhanced_file.path:
                prs = PptxPresentation(presentation.enhanced_file.path)
            elif presentation.original_file and presentation.original_file.path:
                prs = PptxPresentation(presentation.original_file.path)
            else:
                return {
                    'success': False,
                    'error': 'No presentation file found'
                }
            
            current_count = len(prs.slides)
            
            # Convert content_data to SlideContent object
            data = preview.content_data
            
            try:
                # Debug logging
                logger.info(f"Converting preview {preview.id}: {preview.title}")
                if data:
                    logger.info(f"Content blocks count: {len(data.get('content_blocks', []))}")
                else:
                    logger.warning("content_data is None!")
            except Exception as log_error:
                logger.warning(f"Logging failed: {log_error}")
            
            slide_content = SlideContent(
                slide_number=data.get('slide_number', 1) if data else 1,
                title=data.get('title', preview.title) if data else preview.title,
                headline=data.get('headline') if data else None,
                navigation=data.get('navigation') if data else None,
                quote=data.get('quote') if data else None,
                quote_author=data.get('quote_author') if data else None,
                content_blocks=data.get('content_blocks', []) if data else [],
                sources=data.get('sources', []) if data else [],
                visual_notes=data.get('visual_notes', []) if data else []
            )
            
            # Get TemplateCollection if available
            template_collection = presentation.template_collection if hasattr(presentation, 'template_collection') else None
            
            if template_collection:
                logger.info(f"Using TemplateCollection: {template_collection.name}")
            else:
                logger.info("No TemplateCollection available, using fallback analysis")
            
            # Add slide using MarkdownToSlidesHandler
            handler = MarkdownToSlidesHandler(prs, template_collection=template_collection)
            new_slide_idx = handler._create_slide(slide_content)
            
            logger.info(f"Slide created at index {new_slide_idx}, total slides: {len(prs.slides)}")
            
            # Save PPTX with unique name to avoid conflicts
            import os
            from django.core.files import File
            
            # Create temp file with unique name
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                prs.save(tmp_file.name)
                tmp_path = tmp_file.name
            
            # Replace the enhanced file
            with open(tmp_path, 'rb') as f:
                presentation.enhanced_file.save(
                    f'enhanced_{presentation.id}_{len(prs.slides)}.pptx',
                    File(f),
                    save=False
                )
            
            # Clean up temp file
            try:
                os.unlink(tmp_path)
            except:
                pass
            
            # Update preview status
            preview.status = 'converted'
            preview.pptx_slide_number = current_count + 1
            preview.converted_at = timezone.now()
            preview.save()
            
            # Update presentation
            presentation.slide_count_enhanced = len(prs.slides)
            presentation.enhancement_status = 'completed'
            presentation.save()
            
            logger.info(f"Converted preview slide {preview.title} to PPTX position {current_count + 1}")
            
            return {
                'success': True,
                'slide_number': current_count + 1,
                'total_slides': len(prs.slides),
                'preview_id': str(preview.id)
            }
        
        except Exception as e:
            logger.error(f"Error converting preview to PPTX: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e)
            }
    
    def convert_all_previews(self, presentation_id: str) -> Dict:
        """
        Convert all preview slides for a presentation
        
        Args:
            presentation_id: UUID of presentation
            
        Returns:
            Result dict with count of converted slides
        """
        from apps.presentation_studio.models import PreviewSlide
        
        try:
            previews = PreviewSlide.objects.filter(
                presentation__id=presentation_id,
                status='preview'
            ).order_by('preview_order')
            
            converted_count = 0
            failed_count = 0
            errors = []
            
            for preview in previews:
                result = self.convert_preview_to_pptx(str(preview.id))
                if result['success']:
                    converted_count += 1
                else:
                    failed_count += 1
                    errors.append(f"Slide {preview.preview_order}: {result.get('error', 'Unknown error')}")
            
            return {
                'success': failed_count == 0,
                'converted_count': converted_count,
                'failed_count': failed_count,
                'total_previews': len(previews),
                'errors': errors
            }
        
        except Exception as e:
            logger.error(f"Error converting all previews: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e),
                'converted_count': 0,
                'failed_count': 0
            }
    
    def update_preview_content(
        self,
        preview_slide_id: str,
        title: Optional[str] = None,
        content_data: Optional[Dict] = None
    ) -> Dict:
        """
        Update preview slide content
        
        Args:
            preview_slide_id: UUID of preview slide
            title: New title (optional)
            content_data: New content data (optional)
            
        Returns:
            Result dict
        """
        from apps.presentation_studio.models import PreviewSlide
        
        try:
            preview = PreviewSlide.objects.get(id=preview_slide_id)
            
            if preview.status == 'converted':
                return {
                    'success': False,
                    'error': 'Cannot edit converted slide'
                }
            
            if title:
                preview.title = title
                if 'title' in preview.content_data:
                    preview.content_data['title'] = title
            
            if content_data:
                preview.content_data.update(content_data)
            
            preview.save()
            
            return {
                'success': True,
                'preview_id': str(preview.id)
            }
        
        except Exception as e:
            logger.error(f"Error updating preview content: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e)
            }
    
    def delete_preview(self, preview_slide_id: str) -> Dict:
        """
        Delete a preview slide
        
        Args:
            preview_slide_id: UUID of preview slide
            
        Returns:
            Result dict
        """
        from apps.presentation_studio.models import PreviewSlide
        
        try:
            preview = PreviewSlide.objects.get(id=preview_slide_id)
            
            if preview.status == 'converted':
                return {
                    'success': False,
                    'error': 'Cannot delete converted slide'
                }
            
            preview.delete()
            
            return {
                'success': True,
                'message': 'Preview deleted'
            }
        
        except Exception as e:
            logger.error(f"Error deleting preview: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e)
            }
    
    def reorder_previews(self, presentation_id: str, order_map: Dict[str, int]) -> Dict:
        """
        Reorder preview slides
        
        Args:
            presentation_id: UUID of presentation
            order_map: Dict mapping preview_id to new order number
            
        Returns:
            Result dict
        """
        from apps.presentation_studio.models import PreviewSlide
        
        try:
            for preview_id, new_order in order_map.items():
                PreviewSlide.objects.filter(
                    id=preview_id,
                    presentation__id=presentation_id
                ).update(preview_order=new_order)
            
            return {
                'success': True,
                'updated_count': len(order_map)
            }
        
        except Exception as e:
            logger.error(f"Error reordering previews: {e}", exc_info=True)
            return {
                'success': False,
                'error': str(e)
            }
