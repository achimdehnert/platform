"""
Slide Deletion Handler - BaseHandler v2.0
Handles deletion of slides from PPTX presentations
"""

from typing import Dict, Any, List
from pydantic import BaseModel, Field
import logging
from pptx import Presentation as PptxPresentation

from apps.bfagent.handlers.base_handler_v2 import BaseHandler

logger = logging.getLogger(__name__)


class SlideDeletionInput(BaseModel):
    """Input schema for slide deletion"""
    pptx_path: str = Field(..., description="Path to PPTX file")
    slide_numbers: List[int] = Field(..., description="List of slide numbers to delete (1-indexed)")
    validate_before: bool = Field(default=True, description="Validate slide numbers before deletion")


class SlideDeletionOutput(BaseModel):
    """Output schema for slide deletion"""
    success: bool = Field(..., description="Whether deletion was successful")
    slides_deleted: int = Field(..., description="Number of slides deleted")
    slides_remaining: int = Field(..., description="Number of slides remaining")
    deleted_slide_numbers: List[int] = Field(..., description="List of deleted slide numbers")
    errors: List[str] = Field(default_factory=list, description="Any errors encountered")


class SlideDeletionHandler(BaseHandler[SlideDeletionInput, SlideDeletionOutput]):
    """
    Handler for deleting slides from PPTX presentations
    
    Features:
    - Single or multiple slide deletion
    - Validation of slide numbers
    - Preservation of presentation integrity
    - Detailed error reporting
    
    Example:
        handler = SlideDeletionHandler()
        result = handler.execute(
            pptx_path="/path/to/presentation.pptx",
            slide_numbers=[5, 6, 7]
        )
    """
    
    InputSchema = SlideDeletionInput
    OutputSchema = SlideDeletionOutput
    handler_name = "slide_deletion"
    handler_version = "1.0.0"
    domain = "presentation_studio"
    category = "core"
    
    def process(self, validated_input: SlideDeletionInput) -> Dict[str, Any]:
        """
        Process slide deletion
        
        Args:
            validated_input: Validated input containing pptx_path and slide_numbers
            
        Returns:
            Dictionary containing deletion results
        """
        pptx_path = validated_input.pptx_path
        slide_numbers = sorted(validated_input.slide_numbers, reverse=True)  # Delete from end
        validate_before = validated_input.validate_before
        
        errors = []
        deleted_slides = []
        
        try:
            # Load presentation
            prs = PptxPresentation(pptx_path)
            original_count = len(prs.slides)
            
            # Validate slide numbers if requested
            if validate_before:
                validation_errors = self._validate_slide_numbers(slide_numbers, original_count)
                if validation_errors:
                    return {
                        'success': False,
                        'slides_deleted': 0,
                        'slides_remaining': original_count,
                        'deleted_slide_numbers': [],
                        'errors': validation_errors
                    }
            
            # Delete slides (from highest to lowest to maintain indices)
            for slide_num in slide_numbers:
                try:
                    slide_index = slide_num - 1  # Convert to 0-indexed
                    
                    if 0 <= slide_index < len(prs.slides):
                        # Get slide ID and remove
                        slide_id = prs.slides._sldIdLst[slide_index]
                        prs.slides._sldIdLst.remove(slide_id)
                        deleted_slides.append(slide_num)
                        logger.info(f"Deleted slide {slide_num}")
                    else:
                        errors.append(f"Slide {slide_num} out of range")
                        
                except Exception as e:
                    error_msg = f"Failed to delete slide {slide_num}: {str(e)}"
                    errors.append(error_msg)
                    logger.error(error_msg)
            
            # Save modified presentation
            if deleted_slides:
                prs.save(pptx_path)
                logger.info(f"Saved presentation with {len(deleted_slides)} slides deleted")
            
            remaining_count = len(prs.slides)
            
            return {
                'success': len(deleted_slides) > 0,
                'slides_deleted': len(deleted_slides),
                'slides_remaining': remaining_count,
                'deleted_slide_numbers': deleted_slides,
                'errors': errors
            }
            
        except Exception as e:
            error_msg = f"Failed to process slide deletion: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {
                'success': False,
                'slides_deleted': 0,
                'slides_remaining': 0,
                'deleted_slide_numbers': [],
                'errors': [error_msg]
            }
    
    def _validate_slide_numbers(self, slide_numbers: List[int], total_slides: int) -> List[str]:
        """
        Validate slide numbers
        
        Args:
            slide_numbers: List of slide numbers to validate
            total_slides: Total number of slides in presentation
            
        Returns:
            List of validation errors (empty if valid)
        """
        errors = []
        
        if not slide_numbers:
            errors.append("No slide numbers provided")
            return errors
        
        for slide_num in slide_numbers:
            if slide_num < 1:
                errors.append(f"Slide number {slide_num} must be >= 1")
            elif slide_num > total_slides:
                errors.append(f"Slide number {slide_num} exceeds total slides ({total_slides})")
        
        # Check for duplicates
        if len(slide_numbers) != len(set(slide_numbers)):
            errors.append("Duplicate slide numbers provided")
        
        return errors
    
    def delete_slides_after_index(self, pptx_path: str, keep_count: int) -> Dict[str, Any]:
        """
        Helper method to delete all slides after a certain index
        
        Args:
            pptx_path: Path to PPTX file
            keep_count: Number of slides to keep
            
        Returns:
            Deletion results
        """
        try:
            prs = PptxPresentation(pptx_path)
            current_count = len(prs.slides)
            
            if current_count <= keep_count:
                return {
                    'success': True,
                    'slides_deleted': 0,
                    'slides_remaining': current_count,
                    'deleted_slide_numbers': [],
                    'errors': []
                }
            
            # Calculate slides to delete
            slides_to_delete = list(range(keep_count + 1, current_count + 1))
            
            return self.execute(
                pptx_path=pptx_path,
                slide_numbers=slides_to_delete,
                validate_before=False
            )
            
        except Exception as e:
            return {
                'success': False,
                'slides_deleted': 0,
                'slides_remaining': 0,
                'deleted_slide_numbers': [],
                'errors': [str(e)]
            }
