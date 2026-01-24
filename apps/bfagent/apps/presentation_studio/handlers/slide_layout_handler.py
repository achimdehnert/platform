"""
Slide Layout Handler for PPTX Studio
Manages slide layout selection and placeholder identification
"""

import logging
from typing import List, Dict, Optional
from pptx import Presentation as PptxPresentation
from pptx.shapes.placeholder import PlaceholderPicture

logger = logging.getLogger(__name__)


class SlideLayoutHandler:
    """
    Handles slide layout selection and placeholder management
    
    Capabilities:
    - Find appropriate layouts
    - Identify placeholders
    - Match layouts to content types
    - Fallback strategies
    """
    
    def __init__(self):
        # Placeholder type constants
        self.PLACEHOLDER_TITLE = 1
        self.PLACEHOLDER_BODY = 2
        self.PLACEHOLDER_CENTER_TITLE = 3
        self.PLACEHOLDER_SUBTITLE = 4
        self.PLACEHOLDER_DATE = 5
        self.PLACEHOLDER_SLIDE_NUMBER = 6
        self.PLACEHOLDER_FOOTER = 7
        self.PLACEHOLDER_HEADER = 8
        self.PLACEHOLDER_OBJECT = 9
        self.PLACEHOLDER_CHART = 10
        self.PLACEHOLDER_TABLE = 11
        self.PLACEHOLDER_CLIPART = 12
        self.PLACEHOLDER_DIAGRAM = 13
        self.PLACEHOLDER_MEDIA = 14
        self.PLACEHOLDER_PICTURE = 15
    
    def find_layout_with_placeholders(
        self,
        prs: PptxPresentation,
        required_types: List[int],
        preferred_name: Optional[str] = None
    ):
        """
        Find layout with specific placeholder types
        
        Args:
            prs: Presentation object
            required_types: List of required placeholder types
            preferred_name: Optional preferred layout name
            
        Returns:
            Best matching layout or None
        """
        best_match = None
        best_score = 0
        
        for layout in prs.slide_layouts:
            try:
                # Check name match first
                if preferred_name and preferred_name.lower() in layout.name.lower():
                    return layout
                
                # Count matching placeholders
                placeholder_types = self._get_placeholder_types(layout)
                matches = sum(1 for req in required_types if req in placeholder_types)
                
                if matches > best_score:
                    best_score = matches
                    best_match = layout
                    
                # Perfect match - return immediately
                if matches == len(required_types):
                    return layout
                    
            except Exception as e:
                logger.debug(f"Error checking layout {layout.name}: {e}")
                continue
        
        return best_match
    
    def find_title_content_layout(self, prs: PptxPresentation):
        """Find best layout for title + content"""
        return self.find_layout_with_placeholders(
            prs,
            [self.PLACEHOLDER_TITLE, self.PLACEHOLDER_BODY],
            preferred_name="Title and Content"
        )
    
    def find_title_only_layout(self, prs: PptxPresentation):
        """Find best layout for title only"""
        return self.find_layout_with_placeholders(
            prs,
            [self.PLACEHOLDER_TITLE],
            preferred_name="Title Only"
        )
    
    def find_blank_layout(self, prs: PptxPresentation):
        """Find blank layout"""
        for layout in prs.slide_layouts:
            if 'blank' in layout.name.lower():
                return layout
        return prs.slide_layouts[0]  # Fallback to first
    
    def get_layout_by_index(self, prs: PptxPresentation, index: int):
        """Get layout by index with bounds checking"""
        try:
            if 0 <= index < len(prs.slide_layouts):
                return prs.slide_layouts[index]
        except:
            pass
        return prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None
    
    def get_placeholders_info(self, layout) -> List[Dict]:
        """
        Get detailed information about layout placeholders
        
        Args:
            layout: Slide layout
            
        Returns:
            List of placeholder info dicts
        """
        placeholders = []
        
        try:
            for shape in layout.placeholders:
                try:
                    info = {
                        'idx': shape.placeholder_format.idx,
                        'type': shape.placeholder_format.type,
                        'type_name': self._get_type_name(shape.placeholder_format.type),
                        'name': shape.name if hasattr(shape, 'name') else 'Unknown',
                        'has_text_frame': hasattr(shape, 'text_frame')
                    }
                    placeholders.append(info)
                except Exception as e:
                    logger.debug(f"Error reading placeholder: {e}")
                    continue
        except Exception as e:
            logger.error(f"Error getting placeholders: {e}")
        
        return placeholders
    
    def _get_placeholder_types(self, layout) -> List[int]:
        """Get list of placeholder types in layout"""
        types = []
        try:
            for shape in layout.placeholders:
                try:
                    types.append(shape.placeholder_format.type)
                except:
                    continue
        except:
            pass
        return types
    
    def _get_type_name(self, type_id: int) -> str:
        """Get human-readable name for placeholder type"""
        type_names = {
            1: 'Title',
            2: 'Body',
            3: 'Center Title',
            4: 'Subtitle',
            5: 'Date',
            6: 'Slide Number',
            7: 'Footer',
            8: 'Header',
            9: 'Object',
            10: 'Chart',
            11: 'Table',
            12: 'Clip Art',
            13: 'Diagram',
            14: 'Media',
            15: 'Picture'
        }
        return type_names.get(type_id, f'Unknown ({type_id})')
    
    def analyze_presentation_layouts(self, prs: PptxPresentation) -> Dict:
        """
        Analyze all layouts in a presentation
        
        Args:
            prs: Presentation object
            
        Returns:
            Dict with layout analysis
        """
        layouts_info = []
        
        for idx, layout in enumerate(prs.slide_layouts):
            try:
                info = {
                    'index': idx,
                    'name': layout.name,
                    'placeholders': self.get_placeholders_info(layout),
                    'placeholder_count': len(layout.placeholders)
                }
                layouts_info.append(info)
            except Exception as e:
                logger.error(f"Error analyzing layout {idx}: {e}")
                continue
        
        return {
            'total_layouts': len(prs.slide_layouts),
            'layouts': layouts_info
        }
