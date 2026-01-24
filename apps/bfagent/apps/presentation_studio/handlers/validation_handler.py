"""
Validation Handler - BaseHandler v2.0
Validates PPTX presentations for errors and issues
"""

from typing import Dict, Any, List, Optional
from pydantic import BaseModel, Field
from enum import Enum
import logging
from pptx import Presentation as PptxPresentation
from pathlib import Path

from apps.bfagent.handlers.base_handler_v2 import BaseHandler

logger = logging.getLogger(__name__)


class ValidationLevel(str, Enum):
    """Validation severity levels"""
    ERROR = "error"
    WARNING = "warning"
    INFO = "info"


class ValidationIssue(BaseModel):
    """Individual validation issue"""
    level: ValidationLevel
    code: str
    message: str
    slide_number: Optional[int] = None
    details: Optional[Dict[str, Any]] = None


class ValidationInput(BaseModel):
    """Input schema for validation"""
    pptx_path: str = Field(..., description="Path to PPTX file to validate")
    check_file_integrity: bool = Field(default=True, description="Check file can be opened")
    check_slides: bool = Field(default=True, description="Validate slide content")
    check_placeholders: bool = Field(default=True, description="Check placeholders")
    check_images: bool = Field(default=False, description="Validate images")
    max_file_size_mb: int = Field(default=100, description="Maximum file size in MB")


class ValidationOutput(BaseModel):
    """Output schema for validation"""
    is_valid: bool = Field(..., description="Overall validation status")
    issues: List[ValidationIssue] = Field(default_factory=list, description="List of issues found")
    error_count: int = Field(default=0, description="Number of errors")
    warning_count: int = Field(default=0, description="Number of warnings")
    info_count: int = Field(default=0, description="Number of info messages")
    metadata: Dict[str, Any] = Field(default_factory=dict, description="Presentation metadata")


class ValidationHandler(BaseHandler[ValidationInput, ValidationOutput]):
    """
    Handler for validating PPTX presentations
    
    Features:
    - File integrity checks
    - Slide content validation
    - Placeholder validation
    - Image validation
    - Size constraints
    - Detailed issue reporting
    
    Example:
        handler = ValidationHandler()
        result = handler.execute(
            pptx_path="/path/to/presentation.pptx",
            check_slides=True
        )
    """
    
    InputSchema = ValidationInput
    OutputSchema = ValidationOutput
    handler_name = "validation"
    handler_version = "1.0.0"
    domain = "presentation_studio"
    category = "quality"
    
    def process(self, validated_input: ValidationInput) -> Dict[str, Any]:
        """
        Process validation
        
        Args:
            validated_input: Validated input containing validation parameters
            
        Returns:
            Dictionary containing validation results
        """
        pptx_path = validated_input.pptx_path
        issues: List[ValidationIssue] = []
        metadata = {}
        
        try:
            # Check 1: File exists
            if not Path(pptx_path).exists():
                issues.append(ValidationIssue(
                    level=ValidationLevel.ERROR,
                    code="FILE_NOT_FOUND",
                    message=f"File not found: {pptx_path}"
                ))
                return self._build_response(issues, metadata)
            
            # Check 2: File size
            file_size_mb = Path(pptx_path).stat().st_size / (1024 * 1024)
            metadata['file_size_mb'] = round(file_size_mb, 2)
            
            if file_size_mb > validated_input.max_file_size_mb:
                issues.append(ValidationIssue(
                    level=ValidationLevel.WARNING,
                    code="FILE_TOO_LARGE",
                    message=f"File size ({file_size_mb:.2f}MB) exceeds maximum ({validated_input.max_file_size_mb}MB)",
                    details={'file_size_mb': file_size_mb, 'max_size_mb': validated_input.max_file_size_mb}
                ))
            
            # Check 3: File integrity
            if validated_input.check_file_integrity:
                try:
                    prs = PptxPresentation(pptx_path)
                    metadata['total_slides'] = len(prs.slides)
                    metadata['slide_width'] = prs.slide_width
                    metadata['slide_height'] = prs.slide_height
                except Exception as e:
                    issues.append(ValidationIssue(
                        level=ValidationLevel.ERROR,
                        code="CORRUPT_FILE",
                        message=f"Cannot open PPTX file: {str(e)}",
                        details={'error': str(e)}
                    ))
                    return self._build_response(issues, metadata)
            else:
                prs = PptxPresentation(pptx_path)
                metadata['total_slides'] = len(prs.slides)
            
            # Check 4: Validate slides
            if validated_input.check_slides:
                slide_issues = self._validate_slides(prs)
                issues.extend(slide_issues)
            
            # Check 5: Validate placeholders
            if validated_input.check_placeholders:
                placeholder_issues = self._validate_placeholders(prs)
                issues.extend(placeholder_issues)
            
            # Check 6: Validate images
            if validated_input.check_images:
                image_issues = self._validate_images(prs)
                issues.extend(image_issues)
            
            return self._build_response(issues, metadata)
            
        except Exception as e:
            logger.error(f"Validation error: {str(e)}", exc_info=True)
            issues.append(ValidationIssue(
                level=ValidationLevel.ERROR,
                code="VALIDATION_FAILED",
                message=f"Validation failed: {str(e)}",
                details={'error': str(e)}
            ))
            return self._build_response(issues, metadata)
    
    def _validate_slides(self, prs: PptxPresentation) -> List[ValidationIssue]:
        """Validate slide content"""
        issues = []
        
        for idx, slide in enumerate(prs.slides, 1):
            # Check for empty slides
            if not slide.shapes:
                issues.append(ValidationIssue(
                    level=ValidationLevel.WARNING,
                    code="EMPTY_SLIDE",
                    message=f"Slide {idx} is empty (no shapes)",
                    slide_number=idx
                ))
            
            # Check for slides with only title
            text_shapes = [s for s in slide.shapes if hasattr(s, 'text')]
            if len(text_shapes) == 1 and slide.shapes.title:
                issues.append(ValidationIssue(
                    level=ValidationLevel.INFO,
                    code="TITLE_ONLY_SLIDE",
                    message=f"Slide {idx} contains only a title",
                    slide_number=idx
                ))
            
            # Check for very short titles
            if slide.shapes.title and len(slide.shapes.title.text.strip()) < 3:
                issues.append(ValidationIssue(
                    level=ValidationLevel.WARNING,
                    code="SHORT_TITLE",
                    message=f"Slide {idx} has very short or empty title",
                    slide_number=idx,
                    details={'title': slide.shapes.title.text}
                ))
        
        return issues
    
    def _validate_placeholders(self, prs: PptxPresentation) -> List[ValidationIssue]:
        """Validate placeholder usage"""
        issues = []
        
        for idx, slide in enumerate(prs.slides, 1):
            placeholders = [s for s in slide.shapes if hasattr(s, 'placeholder_format')]
            
            for placeholder in placeholders:
                # Check for empty placeholders
                if hasattr(placeholder, 'text') and not placeholder.text.strip():
                    issues.append(ValidationIssue(
                        level=ValidationLevel.INFO,
                        code="EMPTY_PLACEHOLDER",
                        message=f"Slide {idx} has empty placeholder",
                        slide_number=idx,
                        details={'placeholder_type': placeholder.placeholder_format.type}
                    ))
        
        return issues
    
    def _validate_images(self, prs: PptxPresentation) -> List[ValidationIssue]:
        """Validate images in presentation"""
        issues = []
        
        for idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        # Check image size
                        if hasattr(image, 'size'):
                            img_size_mb = image.size / (1024 * 1024)
                            if img_size_mb > 5:  # 5MB threshold
                                issues.append(ValidationIssue(
                                    level=ValidationLevel.WARNING,
                                    code="LARGE_IMAGE",
                                    message=f"Slide {idx} contains large image ({img_size_mb:.2f}MB)",
                                    slide_number=idx,
                                    details={'image_size_mb': img_size_mb}
                                ))
                    except Exception as e:
                        issues.append(ValidationIssue(
                            level=ValidationLevel.ERROR,
                            code="CORRUPT_IMAGE",
                            message=f"Slide {idx} has corrupt image: {str(e)}",
                            slide_number=idx
                        ))
        
        return issues
    
    def _build_response(self, issues: List[ValidationIssue], metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Build validation response"""
        error_count = sum(1 for i in issues if i.level == ValidationLevel.ERROR)
        warning_count = sum(1 for i in issues if i.level == ValidationLevel.WARNING)
        info_count = sum(1 for i in issues if i.level == ValidationLevel.INFO)
        
        return {
            'is_valid': error_count == 0,
            'issues': [i.dict() for i in issues],
            'error_count': error_count,
            'warning_count': warning_count,
            'info_count': info_count,
            'metadata': metadata
        }
