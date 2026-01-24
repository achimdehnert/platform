"""
Django management command to create slides from markdown file
Usage: python manage.py create_slides_from_markdown <markdown_file> [options]
"""

import os
import logging
from django.core.management.base import BaseCommand, CommandError
from pptx import Presentation as PptxPresentation

from apps.presentation_studio.handlers.markdown_slide_parser import parse_markdown_file
from apps.presentation_studio.handlers.markdown_to_slides_handler import (
    MarkdownToSlidesHandler,
    create_presentation_from_markdown
)

logger = logging.getLogger(__name__)


class Command(BaseCommand):
    help = 'Create PowerPoint slides from structured markdown file'

    def add_arguments(self, parser):
        parser.add_argument(
            'markdown_file',
            type=str,
            help='Path to markdown file with slide definitions'
        )
        parser.add_argument(
            '--output',
            '-o',
            type=str,
            help='Output PPTX file path (default: <markdown_file>.pptx)'
        )
        parser.add_argument(
            '--template',
            '-t',
            type=str,
            help='Path to PPTX template file'
        )
        parser.add_argument(
            '--start',
            type=int,
            help='Start slide number (1-indexed)'
        )
        parser.add_argument(
            '--end',
            type=int,
            help='End slide number (1-indexed)'
        )
        parser.add_argument(
            '--summary',
            action='store_true',
            help='Show summary of slides without creating PPTX'
        )

    def handle(self, *args, **options):
        markdown_file = options['markdown_file']
        output_file = options.get('output')
        template_file = options.get('template')
        start_slide = options.get('start')
        end_slide = options.get('end')
        show_summary = options.get('summary')

        # Validate markdown file
        if not os.path.exists(markdown_file):
            raise CommandError(f'Markdown file not found: {markdown_file}')

        # Parse markdown
        self.stdout.write(f'Parsing markdown file: {markdown_file}')
        try:
            parser = parse_markdown_file(markdown_file)
            self.stdout.write(
                self.style.SUCCESS(f'✓ Parsed {len(parser.slides)} slides')
            )
        except Exception as e:
            raise CommandError(f'Error parsing markdown: {e}')

        # Show summary if requested
        if show_summary:
            self.stdout.write('\n' + parser.export_summary())
            return

        # Create presentation
        self.stdout.write('\nCreating PowerPoint presentation...')
        try:
            # Load template or create new
            if template_file:
                if not os.path.exists(template_file):
                    raise CommandError(f'Template file not found: {template_file}')
                prs = PptxPresentation(template_file)
                self.stdout.write(f'  Using template: {template_file}')
            else:
                prs = PptxPresentation()
                self.stdout.write('  Using default template')

            # Create handler
            handler = MarkdownToSlidesHandler(prs)

            # Create slides
            created_slides = handler.create_slides_from_markdown(
                parser,
                start_slide=start_slide,
                end_slide=end_slide
            )

            self.stdout.write(
                self.style.SUCCESS(f'✓ Created {len(created_slides)} slides')
            )

            # Determine output file
            if not output_file:
                base_name = os.path.splitext(markdown_file)[0]
                output_file = f"{base_name}_slides.pptx"

            # Save presentation
            prs.save(output_file)
            self.stdout.write(
                self.style.SUCCESS(f'✓ Saved presentation: {output_file}')
            )

            # Show slide details
            self.stdout.write('\nCreated slides:')
            for idx, slide in enumerate(parser.slides, 1):
                if start_slide and idx < start_slide:
                    continue
                if end_slide and idx > end_slide:
                    break
                self.stdout.write(f'  {idx}. {slide.title}')

        except Exception as e:
            raise CommandError(f'Error creating presentation: {e}')
