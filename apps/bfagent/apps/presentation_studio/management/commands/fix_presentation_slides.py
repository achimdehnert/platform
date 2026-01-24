"""
Management command to fix stuck slides in presentation
"""

from django.core.management.base import BaseCommand
from pptx import Presentation as PptxPresentation
from apps.presentation_studio.models import Presentation


class Command(BaseCommand):
    help = 'Fix presentation with stuck research slides'

    def add_arguments(self, parser):
        parser.add_argument(
            'presentation_id',
            type=str,
            help='Presentation UUID'
        )
        parser.add_argument(
            '--target-count',
            type=int,
            default=13,
            help='Target slide count (default: 13)'
        )

    def handle(self, *args, **options):
        presentation_id = options['presentation_id']
        target_count = options['target_count']
        
        try:
            presentation = Presentation.objects.get(id=presentation_id)
        except Presentation.DoesNotExist:
            self.stdout.write(self.style.ERROR(f'Presentation {presentation_id} not found'))
            return
        
        self.stdout.write(f'Presentation: {presentation.title}')
        self.stdout.write(f'Original file: {presentation.original_file.path}')
        self.stdout.write(f'slide_count_original: {presentation.slide_count_original}')
        
        # Load PPTX
        prs = PptxPresentation(presentation.original_file.path)
        current_slides = len(prs.slides)
        self.stdout.write(f'Current slides in file: {current_slides}')
        
        if current_slides <= target_count:
            self.stdout.write(self.style.SUCCESS(
                f'File already has {current_slides} slides, no fix needed'
            ))
            return
        
        # Remove slides beyond target
        self.stdout.write(self.style.WARNING(
            f'Removing slides {target_count+1} to {current_slides}...'
        ))
        
        while len(prs.slides) > target_count:
            slide_id = prs.slides._sldIdLst[-1]
            prs.slides._sldIdLst.remove(slide_id)
            self.stdout.write(f'  Removed slide, now {len(prs.slides)} slides')
        
        # Save
        prs.save(presentation.original_file.path)
        self.stdout.write(self.style.SUCCESS(
            f'Saved! Now has {len(prs.slides)} slides'
        ))
        
        # Update model
        presentation.slide_count_original = target_count
        presentation.save()
        self.stdout.write(self.style.SUCCESS(
            f'Updated slide_count_original to {target_count}'
        ))
        
        self.stdout.write(self.style.SUCCESS(
            '\n🎉 DONE! Try Research Agent again now!'
        ))
