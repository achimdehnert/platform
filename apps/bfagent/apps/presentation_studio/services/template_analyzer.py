"""
Template Analyzer Service
Analyzes PPTX files to identify reusable slide templates
"""

import logging
from typing import Dict, Optional, List
from pptx import Presentation
from PIL import Image
import io

logger = logging.getLogger(__name__)


class TemplateAnalyzer:
    """
    Analyze PowerPoint presentations to identify reusable slide templates
    
    Classifies slides into template types:
    - title_slide: Title page with minimal content
    - content_slide: Standard content slide with text/images
    - bullet_slide: Slide with bullet points
    - quote_slide: Centered quote or testimonial
    - two_column: Two-column layout
    """
    
    def __init__(self):
        self.slide_types = [
            'title_slide',
            'content_slide',
            'bullet_slide',
            'quote_slide',
            'two_column'
        ]
    
    def analyze_presentation(self, pptx_path: str) -> Dict[str, Dict]:
        """
        Analyze PPTX and return template configuration
        
        Args:
            pptx_path: Path to PPTX file
        
        Returns:
            Dictionary mapping template types to their configurations:
            {
                'title_slide': {
                    'slide_index': 0,
                    'layout_name': 'Title Slide',
                    'shape_count': 2,
                    'has_title': True,
                    'has_content': False
                },
                ...
            }
        """
        try:
            prs = Presentation(pptx_path)
            templates = {}
            
            logger.info(f"Analyzing {len(prs.slides)} slides in {pptx_path}")
            
            for idx, slide in enumerate(prs.slides):
                template_type = self._classify_slide(slide)
                
                if template_type and template_type not in templates:
                    config = self._extract_slide_config(slide, idx)
                    templates[template_type] = config
                    logger.info(f"Identified {template_type} at slide {idx}")
            
            logger.info(f"Found {len(templates)} template types")
            return templates
            
        except Exception as e:
            logger.error(f"Error analyzing presentation: {e}", exc_info=True)
            return {}
    
    def _classify_slide(self, slide) -> Optional[str]:
        """
        Classify slide into template type based on structure
        
        Args:
            slide: python-pptx Slide object
        
        Returns:
            Template type string or None if not classifiable
        """
        shape_count = len(slide.shapes)
        has_title = slide.shapes.title is not None
        
        # Count text shapes and their properties
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        bullet_count = 0
        centered_text = 0
        
        for shape in text_shapes:
            if shape == slide.shapes.title:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.level > 0:
                    bullet_count += 1
                
                # Check if text is centered
                if hasattr(paragraph, 'alignment') and paragraph.alignment == 1:  # CENTER
                    centered_text += 1
        
        # Classification logic
        
        # Title Slide: Only title, maybe subtitle (1-2 shapes)
        if shape_count <= 3 and has_title and len(text_shapes) <= 2:
            return 'title_slide'
        
        # Quote Slide: Centered text, minimal shapes
        if centered_text > 0 and shape_count <= 4:
            return 'quote_slide'
        
        # Bullet Slide: Has bullet points
        if bullet_count >= 2 and has_title:
            return 'bullet_slide'
        
        # Two Column: 4-6 shapes, title + 2 content areas
        if 4 <= shape_count <= 6 and has_title and len(text_shapes) >= 2:
            # Check for side-by-side layout
            if self._has_two_column_layout(slide):
                return 'two_column'
        
        # Content Slide: Title + content (2-4 shapes)
        if 2 <= shape_count <= 5 and has_title and len(text_shapes) >= 1:
            return 'content_slide'
        
        return None
    
    def _has_two_column_layout(self, slide) -> bool:
        """Check if slide has two-column layout"""
        text_shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]
        
        if len(text_shapes) < 2:
            return False
        
        # Check if shapes are side-by-side (x-coordinates differ significantly)
        x_coords = [s.left for s in text_shapes]
        x_coords.sort()
        
        if len(x_coords) >= 2:
            # If there's significant horizontal spacing, it's likely two-column
            return (x_coords[-1] - x_coords[0]) > 3000000  # EMUs
        
        return False
    
    def _extract_slide_config(self, slide, slide_index: int) -> Dict:
        """
        Extract configuration from a slide
        
        Args:
            slide: python-pptx Slide object
            slide_index: Index of slide in presentation
        
        Returns:
            Configuration dictionary
        """
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        image_shapes = [s for s in slide.shapes if hasattr(s, 'image')]
        
        config = {
            'slide_index': slide_index,
            'layout_name': slide.slide_layout.name,
            'shape_count': len(slide.shapes),
            'has_title': slide.shapes.title is not None,
            'has_content': len(text_shapes) > (1 if slide.shapes.title else 0),
            'text_shape_count': len(text_shapes),
            'image_count': len(image_shapes),
        }
        
        # Extract title if present
        if slide.shapes.title:
            config['title_text'] = slide.shapes.title.text[:100]  # First 100 chars
        
        return config
    
    def generate_thumbnail(self, pptx_path: str, slide_index: int, output_path: str, size=(200, 150)):
        """
        Generate thumbnail image for a slide
        
        Args:
            pptx_path: Path to PPTX file
            slide_index: Index of slide to thumbnail
            output_path: Where to save thumbnail
            size: Thumbnail size (width, height)
        
        Returns:
            Path to saved thumbnail or None if failed
        """
        try:
            # This requires additional dependencies (python-pptx doesn't support rendering)
            # Would need aspose.slides or similar for actual thumbnail generation
            # For now, return None as placeholder
            logger.warning("Thumbnail generation not yet implemented")
            return None
            
        except Exception as e:
            logger.error(f"Error generating thumbnail: {e}")
            return None
    
    def validate_templates(self, templates: Dict) -> bool:
        """
        Validate that templates dictionary has required structure
        
        Args:
            templates: Templates dictionary to validate
        
        Returns:
            True if valid, False otherwise
        """
        if not isinstance(templates, dict):
            return False
        
        required_fields = ['slide_index', 'layout_name', 'shape_count']
        
        for template_type, config in templates.items():
            if not isinstance(config, dict):
                return False
            
            for field in required_fields:
                if field not in config:
                    logger.warning(f"Template {template_type} missing required field: {field}")
                    return False
        
        return True
    
    @staticmethod
    def get_supported_template_types() -> List[str]:
        """Get list of supported template types"""
        return [
            'title_slide',
            'content_slide',
            'bullet_slide',
            'quote_slide',
            'two_column'
        ]
