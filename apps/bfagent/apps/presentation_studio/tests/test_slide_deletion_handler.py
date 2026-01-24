"""
Tests for SlideDeletionHandler
"""

import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
from apps.presentation_studio.handlers.slide_deletion_handler import (
    SlideDeletionHandler,
    SlideDeletionInput,
    SlideDeletionOutput
)


class TestSlideDeletionHandler:
    """Test suite for SlideDeletionHandler"""
    
    @pytest.fixture
    def handler(self):
        """Create handler instance"""
        return SlideDeletionHandler()
    
    @pytest.fixture
    def mock_presentation(self):
        """Create mock presentation"""
        mock_prs = Mock()
        mock_prs.slides = Mock()
        mock_prs.slides._sldIdLst = [Mock() for _ in range(10)]  # 10 slides
        mock_prs.slides.__len__ = Mock(return_value=10)
        return mock_prs
    
    def test_should_initialize_handler_correctly(self, handler):
        """Test: Handler should initialize with correct metadata"""
        assert handler.handler_name == "slide_deletion"
        assert handler.handler_version == "1.0.0"
        assert handler.domain == "presentation_studio"
        assert handler.category == "core"
    
    def test_should_validate_input_schema(self, handler):
        """Test: Handler should validate input schema"""
        # Valid input
        valid_input = SlideDeletionInput(
            pptx_path="/test/path.pptx",
            slide_numbers=[1, 2, 3]
        )
        assert valid_input.pptx_path == "/test/path.pptx"
        assert valid_input.slide_numbers == [1, 2, 3]
        
        # Invalid input - no path
        with pytest.raises(Exception):
            SlideDeletionInput(slide_numbers=[1, 2, 3])
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_delete_single_slide_successfully(self, mock_pptx, handler, mock_presentation):
        """Test: Handler should delete a single slide"""
        mock_pptx.return_value = mock_presentation
        
        result = handler.execute(
            pptx_path="/test/path.pptx",
            slide_numbers=[5]
        )
        
        assert result['success'] is True
        assert result['slides_deleted'] == 1
        assert 5 in result['deleted_slide_numbers']
        mock_presentation.save.assert_called_once()
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_delete_multiple_slides_successfully(self, mock_pptx, handler, mock_presentation):
        """Test: Handler should delete multiple slides"""
        mock_pptx.return_value = mock_presentation
        
        result = handler.execute(
            pptx_path="/test/path.pptx",
            slide_numbers=[3, 7, 9]
        )
        
        assert result['success'] is True
        assert result['slides_deleted'] == 3
        assert set(result['deleted_slide_numbers']) == {3, 7, 9}
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_reject_invalid_slide_numbers(self, mock_pptx, handler, mock_presentation):
        """Test: Handler should reject slide numbers out of range"""
        mock_pptx.return_value = mock_presentation
        
        result = handler.execute(
            pptx_path="/test/path.pptx",
            slide_numbers=[0, 15, -1]  # Invalid slide numbers
        )
        
        assert result['success'] is False
        assert len(result['errors']) > 0
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_handle_duplicate_slide_numbers(self, mock_pptx, handler, mock_presentation):
        """Test: Handler should detect duplicate slide numbers"""
        mock_pptx.return_value = mock_presentation
        
        result = handler.execute(
            pptx_path="/test/path.pptx",
            slide_numbers=[3, 3, 5]
        )
        
        assert len([e for e in result['errors'] if 'Duplicate' in e]) > 0
    
    def test_should_validate_slide_numbers_correctly(self, handler):
        """Test: Validation should catch various edge cases"""
        # Empty list
        errors = handler._validate_slide_numbers([], 10)
        assert len(errors) > 0
        
        # Negative numbers
        errors = handler._validate_slide_numbers([-1, 0], 10)
        assert len(errors) > 0
        
        # Out of range
        errors = handler._validate_slide_numbers([15, 20], 10)
        assert len(errors) > 0
        
        # Valid numbers
        errors = handler._validate_slide_numbers([1, 5, 10], 10)
        assert len(errors) == 0
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_delete_slides_in_descending_order(self, mock_pptx, handler, mock_presentation):
        """Test: Handler should delete slides from highest to lowest"""
        mock_pptx.return_value = mock_presentation
        removed_indices = []
        
        def track_removal(slide_id):
            idx = mock_presentation.slides._sldIdLst.index(slide_id)
            removed_indices.append(idx)
            mock_presentation.slides._sldIdLst.remove(slide_id)
        
        mock_presentation.slides._sldIdLst.remove = Mock(side_effect=track_removal)
        
        result = handler.execute(
            pptx_path="/test/path.pptx",
            slide_numbers=[2, 5, 8]
        )
        
        # Should delete in order: 8, 5, 2 (highest to lowest)
        assert removed_indices[0] > removed_indices[1] > removed_indices[2]
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_handle_pptx_load_error(self, mock_pptx, handler):
        """Test: Handler should handle file loading errors"""
        mock_pptx.side_effect = Exception("Cannot open file")
        
        result = handler.execute(
            pptx_path="/test/invalid.pptx",
            slide_numbers=[1]
        )
        
        assert result['success'] is False
        assert result['slides_deleted'] == 0
        assert len(result['errors']) > 0
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_use_delete_slides_after_index_helper(self, mock_pptx, handler, mock_presentation):
        """Test: Helper method should delete all slides after index"""
        mock_pptx.return_value = mock_presentation
        
        result = handler.delete_slides_after_index(
            pptx_path="/test/path.pptx",
            keep_count=5
        )
        
        # Should delete slides 6-10 (5 slides)
        assert result['slides_deleted'] == 5
        assert result['slides_remaining'] == 5
    
    @patch('apps.presentation_studio.handlers.slide_deletion_handler.PptxPresentation')
    def test_should_skip_deletion_when_count_already_correct(self, mock_pptx, handler, mock_presentation):
        """Test: Helper should skip if presentation already has target count"""
        mock_pptx.return_value = mock_presentation
        
        result = handler.delete_slides_after_index(
            pptx_path="/test/path.pptx",
            keep_count=15  # More than current 10
        )
        
        assert result['success'] is True
        assert result['slides_deleted'] == 0
    
    def test_should_format_output_correctly(self, handler):
        """Test: Output should match OutputSchema"""
        output_data = {
            'success': True,
            'slides_deleted': 3,
            'slides_remaining': 7,
            'deleted_slide_numbers': [3, 5, 7],
            'errors': []
        }
        
        output = SlideDeletionOutput(**output_data)
        assert output.success is True
        assert output.slides_deleted == 3
        assert output.slides_remaining == 7


@pytest.mark.integration
class TestSlideDeletionHandlerIntegration:
    """Integration tests with real PPTX files"""
    
    @pytest.fixture
    def sample_pptx_path(self, tmp_path):
        """Create a sample PPTX for testing"""
        from pptx import Presentation
        
        prs = Presentation()
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"Slide {i+1}"
        
        pptx_path = tmp_path / "test_presentation.pptx"
        prs.save(str(pptx_path))
        return str(pptx_path)
    
    def test_should_delete_real_slides(self, sample_pptx_path):
        """Test: Should delete slides from real PPTX file"""
        handler = SlideDeletionHandler()
        
        result = handler.execute(
            pptx_path=sample_pptx_path,
            slide_numbers=[2, 4]
        )
        
        assert result['success'] is True
        assert result['slides_deleted'] == 2
        assert result['slides_remaining'] == 3
        
        # Verify slides were actually deleted
        from pptx import Presentation
        prs = Presentation(sample_pptx_path)
        assert len(prs.slides) == 3
