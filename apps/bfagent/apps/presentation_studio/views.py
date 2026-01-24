"""
Presentation Studio Views
Handle upload, enhancement, and download of PowerPoint presentations
"""

import json
import logging
import time

from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.http import FileResponse, Http404, JsonResponse
from django.shortcuts import get_object_or_404, redirect, render
from django.views.decorators.http import require_http_methods

from apps.core.services.extractors import PDFExtractor, SlideExtractor
from apps.presentation_studio.handlers.enhance_content_handler import EnhanceContentHandler
from apps.presentation_studio.handlers.slide_editor import SlideEditor
from apps.presentation_studio.models import Enhancement, Presentation

logger = logging.getLogger(__name__)


@login_required
def presentation_list(request):
    """List all presentations for current user"""
    presentations = Presentation.objects.filter(uploaded_by=request.user).order_by("-uploaded_at")

    return render(
        request,
        "presentation_studio/presentation_list.html",
        {
            "presentations": presentations,
        },
    )


@login_required
@require_http_methods(["GET", "POST"])
def upload_presentation(request):
    """Upload a new presentation"""
    if request.method == "POST":
        try:
            # Get uploaded file
            pptx_file = request.FILES.get("pptx_file")
            if not pptx_file:
                messages.error(request, "Please select a file to upload")
                return redirect("presentation_studio:list")

            # Validate file type
            if not (pptx_file.name.endswith(".pptx") or pptx_file.name.endswith(".json")):
                messages.error(request, "Only .pptx and .json files are allowed")
                return redirect("presentation_studio:list")

            # Get metadata first
            title = request.POST.get("title", pptx_file.name)
            description = request.POST.get("description", "")

            # Handle JSON upload (convert to PPTX first)
            if pptx_file.name.endswith(".json"):
                try:
                    import json
                    import tempfile

                    from pptx import Presentation as PptxPresentation

                    from apps.presentation_studio.handlers.json_to_slides_handler import (
                        JsonToSlidesHandler,
                    )

                    # Parse JSON
                    json_content = json.load(pptx_file)

                    # Create PPTX from JSON
                    prs = PptxPresentation()
                    handler = JsonToSlidesHandler(prs)
                    handler.create_slides_from_json(json_content)

                    # Save to temp file
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                        prs.save(tmp.name)
                        tmp_path = tmp.name

                    # Replace pptx_file with generated PPTX
                    from django.core.files import File

                    with open(tmp_path, "rb") as f:
                        pptx_file = File(f, name=pptx_file.name.replace(".json", ".pptx"))

                        # Create presentation
                        presentation = Presentation.objects.create(
                            title=title,
                            description=description + " (created from JSON)",
                            original_file=pptx_file,
                            uploaded_by=request.user,
                            enhancement_status="uploaded",
                        )

                    # Clean up
                    import os

                    try:
                        os.unlink(tmp_path)
                    except:
                        pass

                    messages.success(
                        request, f'Presentation "{title}" created from JSON successfully'
                    )
                    return redirect("presentation_studio:detail", pk=presentation.id)

                except Exception as e:
                    logger.error(f"JSON conversion error: {str(e)}", exc_info=True)
                    messages.error(request, f"JSON conversion failed: {str(e)}")
                    return redirect("presentation_studio:list")

            # Get metadata
            title = request.POST.get("title", pptx_file.name)
            description = request.POST.get("description", "")

            # Create presentation
            presentation = Presentation.objects.create(
                title=title,
                description=description,
                original_file=pptx_file,
                uploaded_by=request.user,
                enhancement_status="uploaded",
            )

            # Extract slide count
            try:
                handler = EnhanceContentHandler()
                slide_count = handler.extract_slide_count(presentation.original_file.path)
                presentation.slide_count_original = slide_count
                presentation.enhancement_status = "ready"
                presentation.save()
            except Exception as e:
                logger.warning(f"Could not extract slide count: {str(e)}")
                presentation.slide_count_original = 0
                presentation.save()

            # Auto-create TemplateCollection from uploaded PPTX
            try:
                from apps.presentation_studio.models import TemplateCollection
                from apps.presentation_studio.services.template_analyzer import TemplateAnalyzer

                # Check if default collection exists
                default_collection = TemplateCollection.get_default()

                if not default_collection:
                    # No default exists, create one from this PPTX
                    logger.info(f"Auto-creating TemplateCollection from {presentation.title}")

                    analyzer = TemplateAnalyzer()
                    templates = analyzer.analyze_presentation(presentation.original_file.path)

                    if templates:
                        import shutil
                        import tempfile

                        from django.core.files import File

                        # Create temp copy of PPTX for collection
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                            shutil.copy2(presentation.original_file.path, tmp.name)
                            tmp_path = tmp.name

                        # Create collection
                        with open(tmp_path, "rb") as f:
                            collection = TemplateCollection.objects.create(
                                name=f"Auto: {presentation.title}",
                                description=f"Auto-generated from uploaded presentation",
                                templates=templates,
                                created_by=request.user,
                                is_default=True,
                                is_active=True,
                            )
                            collection.master_pptx.save(
                                f"auto_{presentation.id}.pptx", File(f), save=True
                            )

                        # Link presentation to collection
                        presentation.template_collection = collection
                        presentation.save()

                        logger.info(
                            f"✓ Created TemplateCollection: {collection.name} ({collection.template_count} templates)"
                        )

                        # Clean up temp file
                        import os

                        try:
                            os.unlink(tmp_path)
                        except:
                            pass
                else:
                    # Default exists, just link to it
                    presentation.template_collection = default_collection
                    presentation.save()
                    logger.info(
                        f"✓ Linked to existing default TemplateCollection: {default_collection.name}"
                    )

            except Exception as e:
                logger.warning(f"Could not auto-create TemplateCollection: {str(e)}")
                # Don't fail the upload if template creation fails

            messages.success(request, f'Presentation "{title}" uploaded successfully')
            return redirect("presentation_studio:detail", pk=presentation.id)

        except Exception as e:
            logger.error(f"Upload error: {str(e)}", exc_info=True)
            messages.error(request, f"Upload failed: {str(e)}")
            return redirect("presentation_studio:list")

    # GET request - show upload form
    return render(request, "presentation_studio/upload.html")


@login_required
def presentation_detail(request, pk):
    """View presentation details"""
    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    # Get enhancement history
    enhancements = presentation.enhancements.all().order_by("-executed_at")

    # Get preview slides
    from apps.presentation_studio.models import PreviewSlide

    preview_slides = PreviewSlide.objects.filter(presentation=presentation).order_by(
        "preview_order"
    )

    # Extract table of contents (prefer enhanced file if available)
    table_of_contents = []

    # Use enhanced file if available, otherwise use original
    if presentation.enhanced_file:
        try:
            handler = EnhanceContentHandler()
            file_path = presentation.enhanced_file.path
            table_of_contents = handler.extract_table_of_contents(file_path)
        except Exception as e:
            logger.warning(f"Could not extract TOC from enhanced file: {str(e)}")
    elif presentation.original_file:
        try:
            handler = EnhanceContentHandler()
            file_path = presentation.original_file.path
            table_of_contents = handler.extract_table_of_contents(file_path)
        except Exception as e:
            logger.warning(f"Could not extract TOC from original file: {str(e)}")

    return render(
        request,
        "presentation_studio/presentation_detail.html",
        {
            "presentation": presentation,
            "enhancements": enhancements,
            "table_of_contents": table_of_contents,
            "preview_slides": preview_slides,
        },
    )


@login_required
@require_http_methods(["POST"])
def enhance_presentation(request, pk):
    """Enhance a presentation with new content"""
    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    try:
        # Get enhancement parameters
        enhancement_source = request.POST.get("enhancement_source", "template")
        enhancement_type = request.POST.get("enhancement_type", "technical")
        enhancement_mode = request.POST.get("enhancement_mode", "append")

        # Handle concepts from file or template
        concepts = []

        if enhancement_source == "file" and request.FILES.get("concepts_file"):
            # Read concepts from uploaded file
            concepts_file = request.FILES["concepts_file"]
            file_name = concepts_file.name.lower()

            try:
                # Check file type
                if file_name.endswith(".pdf"):
                    # PDF file - extract content
                    import os
                    import tempfile

                    # Save PDF temporarily
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                        for chunk in concepts_file.chunks():
                            tmp_file.write(chunk)
                        tmp_path = tmp_file.name

                    # Extract content from PDF
                    pdf_extractor = PDFExtractor()
                    concepts = pdf_extractor.get_available_concepts(tmp_path)

                    # Clean up temp file
                    try:
                        os.unlink(tmp_path)
                    except:
                        pass

                    if not concepts:
                        return JsonResponse(
                            {
                                "success": False,
                                "error": "Could not extract content from PDF. Make sure pdfplumber is installed.",
                            },
                            status=400,
                        )

                elif file_name.endswith(".json") or file_name.endswith(".txt"):
                    # JSON file - parse directly
                    file_content = concepts_file.read().decode("utf-8")
                    concepts = json.loads(file_content)

                    # Validate concepts format
                    if not isinstance(concepts, list):
                        return JsonResponse(
                            {"success": False, "error": "Concepts file must contain a JSON array"},
                            status=400,
                        )

                    # Validate each concept has title
                    for concept in concepts:
                        if "title" not in concept:
                            return JsonResponse(
                                {
                                    "success": False,
                                    "error": 'Each concept must have a "title" field',
                                },
                                status=400,
                            )
                        # Add default content if not provided
                        if "content" not in concept:
                            concept["content"] = ""

                elif file_name.endswith(".md") or file_name.endswith(".markdown"):
                    # Markdown file - create preview slides instead of direct conversion
                    import os
                    import tempfile

                    from apps.presentation_studio.handlers.preview_slide_handler import (
                        PreviewSlideHandler,
                    )

                    tmp_path = None
                    try:
                        # Save markdown temporarily
                        with tempfile.NamedTemporaryFile(
                            delete=False, suffix=".md", mode="w", encoding="utf-8"
                        ) as tmp_file:
                            file_content = concepts_file.read().decode("utf-8")
                            tmp_file.write(file_content)
                            tmp_path = tmp_file.name

                        logger.info(
                            f"Processing markdown file: {concepts_file.name}, temp: {tmp_path}"
                        )

                        # Create preview slides
                        handler = PreviewSlideHandler()
                        preview_slides = handler.create_previews_from_markdown(
                            presentation, tmp_path, file_name=concepts_file.name
                        )

                        logger.info(
                            f"Created {len(preview_slides) if preview_slides else 0} preview slides"
                        )

                        if not preview_slides:
                            return JsonResponse(
                                {
                                    "success": False,
                                    "error": "No slides found in markdown file. Check format.",
                                },
                                status=400,
                            )

                        # Return success with preview info (no PPTX generation yet)
                        return JsonResponse(
                            {
                                "success": True,
                                "message": f'Created {len(preview_slides)} preview slides. Use "Convert to PPTX" buttons to add them to your presentation.',
                                "preview_mode": True,
                                "preview_count": len(preview_slides),
                                "slides_before": presentation.slide_count_enhanced
                                or presentation.slide_count_original
                                or 0,
                            }
                        )

                    except Exception as markdown_error:
                        logger.error(
                            f"Markdown processing error: {str(markdown_error)}", exc_info=True
                        )
                        return JsonResponse(
                            {
                                "success": False,
                                "error": f"Error processing markdown file: {str(markdown_error)}",
                            },
                            status=400,
                        )

                    finally:
                        # Clean up temp file
                        if tmp_path:
                            try:
                                os.unlink(tmp_path)
                                logger.info(f"Cleaned up temp file: {tmp_path}")
                            except Exception as cleanup_error:
                                logger.warning(
                                    f"Could not delete temp file {tmp_path}: {cleanup_error}"
                                )

                else:
                    return JsonResponse(
                        {
                            "success": False,
                            "error": "Unsupported file type. Please upload PDF, JSON, or Markdown file.",
                        },
                        status=400,
                    )

            except json.JSONDecodeError as e:
                return JsonResponse(
                    {"success": False, "error": f"Invalid JSON format in file: {str(e)}"},
                    status=400,
                )
            except Exception as e:
                return JsonResponse(
                    {"success": False, "error": f"Error reading file: {str(e)}"}, status=400
                )
        else:
            # Use template-based concepts
            concepts_json = request.POST.get("concepts", "[]")

            try:
                concepts = json.loads(concepts_json)
            except json.JSONDecodeError:
                return JsonResponse(
                    {"success": False, "error": "Invalid concepts format"}, status=400
                )

            # If no concepts provided, use default templates
            if not concepts:
                handler = EnhanceContentHandler()
                templates = handler.get_available_concepts(enhancement_type)
                concepts = [
                    {"title": t["title"], "content": "Generated content"} for t in templates
                ]

        # Update status
        presentation.enhancement_status = "enhancing"
        presentation.save()

        # Track start time
        start_time = time.time()

        # Call Enhancement Handler
        handler = EnhanceContentHandler()
        result = handler.enhance_presentation(
            original_path=presentation.original_file.path,
            enhancement_type=enhancement_type,
            concepts=concepts,
            mode=enhancement_mode,
        )

        # Calculate duration
        duration = time.time() - start_time

        # Create enhancement record
        enhancement = Enhancement.objects.create(
            presentation=presentation,
            enhancement_type=enhancement_type,
            enhancement_mode=enhancement_mode,
            concepts=concepts,
            slides_before=result.get("slides_before", 0),
            slides_after=result.get("slides_after", 0),
            executed_by=request.user,
            success=result.get("success", False),
            error_message=result.get("error", ""),
            result_data=result,
            duration_seconds=duration,
        )

        # Update presentation
        if result["success"]:
            # Read enhanced file and save to Django FileField
            with open(result["enhanced_path"], "rb") as f:
                from django.core.files import File

                presentation.enhanced_file.save(
                    f"enhanced_{presentation.id}.pptx", File(f), save=False
                )

            presentation.slide_count_enhanced = result["slides_after"]
            presentation.enhancement_status = "completed"
            presentation.concepts_added = concepts
            presentation.save()

            messages.success(
                request, f"Enhancement completed! Added {result['slides_added']} slides."
            )
            return JsonResponse(
                {
                    "success": True,
                    "enhancement_id": str(enhancement.id),
                    "slides_added": result["slides_added"],
                    "message": f"Successfully enhanced presentation with {result['slides_added']} new slides",
                }
            )
        else:
            presentation.enhancement_status = "failed"
            presentation.save()

            return JsonResponse(
                {
                    "success": False,
                    "error": result.get("error", "Enhancement failed"),
                    "enhancement_id": str(enhancement.id),
                },
                status=500,
            )

    except Exception as e:
        logger.error(f"Enhancement error: {str(e)}", exc_info=True)
        presentation.enhancement_status = "failed"
        presentation.save()

        return JsonResponse({"success": False, "error": str(e)}, status=500)


@login_required
def download_enhanced(request, pk):
    """Download enhanced presentation"""
    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    # Check if enhanced file exists
    if not presentation.enhanced_file:
        raise Http404("Enhanced file not available")

    # Return file
    response = FileResponse(
        presentation.enhanced_file.open("rb"),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    response["Content-Disposition"] = f'attachment; filename="{presentation.title}_enhanced.pptx"'

    return response


@login_required
def slide_viewer(request, pk, slide_number):
    """View detailed slide information using medtrans XML extractor"""
    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    try:
        # Use new SlideExtractor for shape-based extraction
        extractor = SlideExtractor()

        # Use enhanced file if available, otherwise use original
        if presentation.enhanced_file:
            file_path = presentation.enhanced_file.path
        else:
            file_path = presentation.original_file.path

        # Extract all slides with shapes
        extraction_result = extractor.extract_slides_with_shapes(file_path)

        if not extraction_result["success"]:
            messages.error(request, "Failed to extract slide data")
            return redirect("presentation_studio:detail", pk=pk)

        # Find the requested slide
        slide_data = None
        for slide in extraction_result["slides"]:
            if slide["slide_number"] == slide_number:
                slide_data = slide
                break

        if not slide_data:
            messages.error(request, f"Slide {slide_number} not found")
            return redirect("presentation_studio:detail", pk=pk)

        # Add navigation info
        total_slides = len(extraction_result["slides"])
        has_prev = slide_number > 1
        has_next = slide_number < total_slides

        context = {
            "presentation": presentation,
            "slide": slide_data,
            "slide_number": slide_number,
            "total_slides": total_slides,
            "has_prev": has_prev,
            "has_next": has_next,
            "prev_slide": slide_number - 1 if has_prev else None,
            "next_slide": slide_number + 1 if has_next else None,
        }

        return render(request, "presentation_studio/slide_viewer.html", context)

    except Exception as e:
        logger.error(f"Slide viewer error: {str(e)}", exc_info=True)
        messages.error(request, f"Error viewing slide: {str(e)}")
        return redirect("presentation_studio:detail", pk=pk)


@login_required
@require_http_methods(["POST"])
def edit_slide(request, pk, slide_number):
    """Edit slide text content"""
    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    try:
        # Parse text updates from request
        text_updates = json.loads(request.POST.get("text_updates", "[]"))

        if not text_updates:
            return JsonResponse({"success": False, "error": "No updates provided"})

        # Create output path for edited version
        import os

        from django.core.files.base import ContentFile

        # Use enhanced file if available, otherwise use original
        if presentation.enhanced_file:
            source_path = presentation.enhanced_file.path
        else:
            source_path = presentation.original_file.path

        temp_output = source_path.replace(".pptx", "_edited_temp.pptx")

        # Use SlideEditor to update texts
        editor = SlideEditor()
        result = editor.update_slide_texts(
            pptx_path=source_path,
            output_path=temp_output,
            slide_number=slide_number,
            text_updates=text_updates,
        )

        if result["success"]:
            # Save edited version (update enhanced file if exists, otherwise original)
            with open(temp_output, "rb") as f:
                content = f.read()
                if presentation.enhanced_file:
                    presentation.enhanced_file.save(
                        presentation.enhanced_file.name, ContentFile(content), save=True
                    )
                else:
                    presentation.original_file.save(
                        presentation.original_file.name, ContentFile(content), save=True
                    )

            # Clean up temp file
            if os.path.exists(temp_output):
                os.remove(temp_output)

            messages.success(request, f'Updated {result["updated_count"]} texts successfully')
            return JsonResponse({"success": True, "updated_count": result["updated_count"]})
        else:
            return JsonResponse({"success": False, "error": ", ".join(result["errors"])})

    except Exception as e:
        logger.error(f"Edit slide error: {str(e)}", exc_info=True)
        return JsonResponse({"success": False, "error": str(e)}, status=500)


@login_required
@require_http_methods(["POST"])
def delete_slide(request, pk, slide_number):
    """
    Delete a specific slide from the presentation
    """
    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    try:
        from pptx import Presentation as PptxPresentation

        # Determine which file to modify
        if presentation.enhanced_file:
            pptx_path = presentation.enhanced_file.path
        else:
            pptx_path = presentation.original_file.path

        # Load presentation
        prs = PptxPresentation(pptx_path)
        total_slides = len(prs.slides)

        # Validate slide number
        if slide_number < 1 or slide_number > total_slides:
            return JsonResponse(
                {
                    "success": False,
                    "error": f"Invalid slide number. Must be between 1 and {total_slides}",
                },
                status=400,
            )

        # Delete slide (convert to 0-indexed)
        slide_index = slide_number - 1

        # Get the slide's rId
        slide_id = prs.slides._sldIdLst[slide_index]

        # Remove from slides collection
        prs.slides._sldIdLst.remove(slide_id)

        # Save modified presentation
        prs.save(pptx_path)

        logger.info(
            f"Deleted slide {slide_number} from presentation {pk}. "
            f"Total slides: {total_slides} -> {total_slides-1}"
        )

        messages.success(request, f"Slide {slide_number} deleted successfully!")

        return JsonResponse(
            {
                "success": True,
                "message": f"Slide {slide_number} deleted successfully",
                "slides_remaining": total_slides - 1,
            }
        )

    except Exception as e:
        logger.error(f"Delete slide error: {str(e)}", exc_info=True)
        return JsonResponse(
            {"success": False, "error": f"Failed to delete slide: {str(e)}"}, status=500
        )


@login_required
@require_http_methods(["POST"])
def delete_presentation(request, pk):
    """
    Delete an entire presentation using PresentationDeletionHandler
    """
    from apps.presentation_studio.handlers.presentation_deletion_handler import (
        PresentationDeletionHandler,
    )

    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    try:
        # Store info for message
        title = presentation.title or "Untitled"
        presentation_id = str(pk)

        # Get file paths
        original_path = presentation.original_file.path if presentation.original_file else None
        enhanced_path = presentation.enhanced_file.path if presentation.enhanced_file else None

        # Use handler to delete files
        handler = PresentationDeletionHandler()
        result = handler.execute(
            {
                "presentation_id": presentation_id,
                "delete_files": True,
                "original_file_path": original_path,
                "enhanced_file_path": enhanced_path,
            }
        )

        # Delete database entry (cascades to enhancements, previews, etc.)
        presentation.delete()

        logger.info(
            f"Deleted presentation {pk} ('{title}') by user {request.user.username}. "
            f"Files deleted: {result['files_deleted']}"
        )

        # Log any file deletion errors
        if result["errors"]:
            for error in result["errors"]:
                logger.warning(f"File deletion warning: {error}")

        messages.success(request, f'Presentation "{title}" deleted successfully!')

        return JsonResponse(
            {
                "success": True,
                "message": f'Presentation "{title}" deleted successfully',
                "files_deleted": result["files_deleted"],
            }
        )

    except Exception as e:
        logger.error(f"Delete presentation error: {str(e)}", exc_info=True)
        return JsonResponse(
            {"success": False, "error": f"Failed to delete presentation: {str(e)}"}, status=500
        )


@login_required
@require_http_methods(["POST"])
def convert_preview_slide(request, pk, preview_pk):
    """Convert a single preview slide to PPTX"""
    from apps.presentation_studio.handlers.preview_slide_handler import PreviewSlideHandler
    from apps.presentation_studio.models import Presentation, PreviewSlide

    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    preview = get_object_or_404(PreviewSlide, pk=preview_pk, presentation=presentation)

    if preview.status == "converted":
        return JsonResponse({"success": False, "error": "Slide already converted"})

    handler = PreviewSlideHandler()
    result = handler.convert_preview_to_pptx(str(preview.pk))

    return JsonResponse(result)


@login_required
@require_http_methods(["POST"])
def convert_all_previews(request, pk):
    """Convert all preview slides for a presentation"""
    from apps.presentation_studio.handlers.preview_slide_handler import PreviewSlideHandler
    from apps.presentation_studio.models import Presentation

    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    handler = PreviewSlideHandler()
    result = handler.convert_all_previews(str(presentation.pk))

    return JsonResponse(result)


@login_required
def get_preview_slides(request, pk):
    """Get all preview slides for a presentation"""
    from apps.presentation_studio.models import Presentation, PreviewSlide

    presentation = get_object_or_404(Presentation, pk=pk, uploaded_by=request.user)

    previews = PreviewSlide.objects.filter(presentation=presentation).order_by("preview_order")

    preview_data = []
    for preview in previews:
        preview_data.append(
            {
                "id": str(preview.id),
                "title": preview.title,
                "preview_order": preview.preview_order,
                "status": preview.status,
                "source_type": preview.source_type,
                "created_at": preview.created_at.isoformat(),
                "pptx_slide_number": preview.pptx_slide_number,
                "can_convert": preview.can_convert,
            }
        )

    return JsonResponse(
        {
            "success": True,
            "previews": preview_data,
            "total_count": len(preview_data),
            "convertible_count": len([p for p in previews if p.can_convert]),
        }
    )
