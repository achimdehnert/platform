"""
Research Agent Views for PPTX Studio
"""

import logging
from django.shortcuts import render, get_object_or_404, redirect
from django.http import JsonResponse
from django.contrib.auth.decorators import login_required
from django.views.decorators.http import require_http_methods
from django.contrib import messages
import json
import os

from apps.presentation_studio.models import Presentation, Enhancement
from apps.presentation_studio.handlers.research_agent_handler import ResearchAgentHandler
from apps.presentation_studio.handlers.enhance_content_handler import EnhanceContentHandler

logger = logging.getLogger(__name__)


@login_required
def research_interface(request, pk):
    """
    Research interface for a presentation
    """
    presentation = get_object_or_404(
        Presentation,
        pk=pk,
        uploaded_by=request.user
    )
    
    return render(request, 'presentation_studio/research_interface.html', {
        'presentation': presentation,
    })


@login_required
@require_http_methods(["POST"])
def perform_research(request, pk):
    """
    Perform research based on prompt
    """
    presentation = get_object_or_404(
        Presentation,
        pk=pk,
        uploaded_by=request.user
    )
    
    try:
        # Get research prompt
        prompt = request.POST.get('research_prompt', '').strip()
        
        if not prompt:
            return JsonResponse({
                'success': False,
                'error': 'Please provide a research prompt'
            }, status=400)
        
        # Optional parameters
        max_results = int(request.POST.get('max_results', 10))
        sources = request.POST.getlist('sources[]') or ['web', 'knowledge']
        
        # Perform research
        research_handler = ResearchAgentHandler()
        result = research_handler.research(
            prompt=prompt,
            options={
                'max_results': max_results,
                'sources': sources
            }
        )
        
        if not result['success']:
            return JsonResponse({
                'success': False,
                'error': result.get('error', 'Research failed')
            }, status=500)
        
        # Store research results in session for later use
        request.session[f'research_results_{pk}'] = {
            'prompt': prompt,
            'findings': result['key_findings'],
            'slide_concepts': result['slide_concepts']
        }
        
        return JsonResponse({
            'success': True,
            'prompt': prompt,
            'total_results': result['total_results'],
            'findings': result['key_findings'],
            'slide_concepts_count': len(result['slide_concepts'])
        })
        
    except ValueError as e:
        return JsonResponse({
            'success': False,
            'error': f'Invalid input: {str(e)}'
        }, status=400)
    except Exception as e:
        logger.error(f'Research error: {str(e)}', exc_info=True)
        return JsonResponse({
            'success': False,
            'error': f'Research failed: {str(e)}'
        }, status=500)


@login_required
@require_http_methods(["POST"])
def generate_slides_from_research(request, pk):
    """
    Generate slides from selected research findings
    """
    presentation = get_object_or_404(
        Presentation,
        pk=pk,
        uploaded_by=request.user
    )
    
    try:
        # Get research results from session
        session_key = f'research_results_{pk}'
        research_data = request.session.get(session_key)
        
        if not research_data:
            return JsonResponse({
                'success': False,
                'error': 'No research data found. Please perform research first.'
            }, status=400)
        
        # Get selected finding IDs
        selected_ids_str = request.POST.get('selected_findings', '[]')
        selected_ids = json.loads(selected_ids_str)
        
        # Filter slide concepts based on selection
        research_handler = ResearchAgentHandler()
        all_concepts = research_data['slide_concepts']
        filtered_concepts = research_handler.filter_concepts(
            all_concepts,
            selected_ids
        )
        
        if not filtered_concepts:
            return JsonResponse({
                'success': False,
                'error': 'No slides selected'
            }, status=400)
        
        # Use EnhanceContentHandler to generate slides
        enhancement_handler = EnhanceContentHandler()
        
        # Get original file path
        if presentation.enhanced_file:
            source_path = presentation.enhanced_file.path
        else:
            source_path = presentation.original_file.path
        
        # Generate enhanced presentation
        import time
        start_time = time.time()
        
        result = enhancement_handler.enhance_presentation(
            original_path=source_path,
            enhancement_type='research',
            concepts=filtered_concepts,
            mode='append'
        )
        
        duration = time.time() - start_time
        
        if not result['success']:
            return JsonResponse({
                'success': False,
                'error': result.get('error', 'Enhancement failed')
            }, status=500)
        
        # Save enhancement record
        Enhancement.objects.create(
            presentation=presentation,
            enhancement_type='research',
            enhancement_mode='append',
            concepts=filtered_concepts,
            slides_before=result['slides_before'],
            slides_after=result['slides_after'],
            success=True,
            result_data=result,
            executed_by=request.user,
            duration_seconds=duration
        )
        
        # Update presentation
        from django.core.files import File
        with open(result['enhanced_path'], 'rb') as f:
            presentation.enhanced_file.save(
                f"enhanced_{presentation.id}.pptx",
                File(f),
                save=False
            )
        
        presentation.slide_count_enhanced = result['slides_after']
        presentation.enhancement_status = 'completed'
        presentation.concepts_added = filtered_concepts
        presentation.save()
        
        # Clear research data from session
        if session_key in request.session:
            del request.session[session_key]
        
        messages.success(
            request,
            f"Generated {result['slides_added']} slides from research!"
        )
        
        return JsonResponse({
            'success': True,
            'slides_added': result['slides_added'],
            'slides_total': result['slides_after'],
            'message': f"Successfully added {result['slides_added']} research slides"
        })
        
    except json.JSONDecodeError:
        return JsonResponse({
            'success': False,
            'error': 'Invalid selection data'
        }, status=400)
    except Exception as e:
        logger.error(f'Slide generation error: {str(e)}', exc_info=True)
        return JsonResponse({
            'success': False,
            'error': f'Failed to generate slides: {str(e)}'
        }, status=500)


@login_required
@require_http_methods(["POST"])
def clear_research_slides(request, pk):
    """
    Reset presentation to original state (remove enhanced version)
    This allows starting fresh with new research
    """
    presentation = get_object_or_404(
        Presentation,
        pk=pk,
        uploaded_by=request.user
    )
    
    try:
        from pptx import Presentation as PptxPresentation
        
        # Strategy: Remove all slides beyond original count from BOTH files
        original_count = presentation.slide_count_original or 13  # Fallback
        
        # Clean enhanced file if exists
        if presentation.enhanced_file and os.path.exists(presentation.enhanced_file.path):
            prs = PptxPresentation(presentation.enhanced_file.path)
            current_count = len(prs.slides)
            
            # Remove slides beyond original count
            while len(prs.slides) > original_count:
                slide_id = prs.slides._sldIdLst[-1]  # Last slide
                prs.slides._sldIdLst.remove(slide_id)
            
            prs.save(presentation.enhanced_file.path)
            logger.info(f"Cleaned enhanced file: {current_count} -> {len(prs.slides)} slides")
            
            # Delete enhanced file completely
            os.remove(presentation.enhanced_file.path)
            presentation.enhanced_file = None
            presentation.save()
        
        # Also clean original file if it has extra slides
        if os.path.exists(presentation.original_file.path):
            prs = PptxPresentation(presentation.original_file.path)
            current_count = len(prs.slides)
            
            if current_count > original_count:
                # Remove slides beyond original count
                while len(prs.slides) > original_count:
                    slide_id = prs.slides._sldIdLst[-1]  # Last slide
                    prs.slides._sldIdLst.remove(slide_id)
                
                prs.save(presentation.original_file.path)
                logger.info(f"Cleaned original file: {current_count} -> {len(prs.slides)} slides")
        
        # Clear session research data
        session_key = f'research_results_{pk}'
        if session_key in request.session:
            del request.session[session_key]
        
        messages.success(
            request,
            f"Research slides cleared! Reset to original {original_count} slides."
        )
        
        return JsonResponse({
            'success': True,
            'message': f'Research slides cleared. Reset to {original_count} slides.'
        })
        
    except Exception as e:
        logger.error(f'Clear research slides error: {str(e)}', exc_info=True)
        return JsonResponse({
            'success': False,
            'error': f'Failed to clear research slides: {str(e)}'
        }, status=500)
