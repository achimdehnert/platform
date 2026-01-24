"""
Emergency fix script for presentation with stuck slides
"""

from pptx import Presentation as PptxPresentation
from apps.presentation_studio.models import Presentation

# Get presentation
presentation_id = 'bb60cc50-b247-4965-9b23-5ee55c722115'
presentation = Presentation.objects.get(id=presentation_id)

print(f"Presentation: {presentation.title}")
print(f"Original file: {presentation.original_file.path}")
print(f"slide_count_original: {presentation.slide_count_original}")

# Load PPTX
prs = PptxPresentation(presentation.original_file.path)
current_slides = len(prs.slides)
print(f"Current slides in file: {current_slides}")

# FORCEFULLY remove slides beyond 13
target_count = 13
if current_slides > target_count:
    print(f"\nRemoving slides {target_count+1} to {current_slides}...")
    
    while len(prs.slides) > target_count:
        slide_id = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(slide_id)
        print(f"  Removed slide, now {len(prs.slides)} slides")
    
    # Save
    prs.save(presentation.original_file.path)
    print(f"\n✅ Saved! Now has {len(prs.slides)} slides")
    
    # Update model
    presentation.slide_count_original = target_count
    presentation.save()
    print(f"✅ Updated slide_count_original to {target_count}")
else:
    print(f"\n✅ File already has {current_slides} slides, no fix needed")

print("\n🎉 DONE! Try Research Agent again now!")
