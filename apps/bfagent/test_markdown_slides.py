"""
Quick test script for markdown to slides conversion
"""

import os
import sys

# Add project to path
sys.path.insert(0, os.path.dirname(__file__))

# Setup Django
os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'config.settings.development')
import django
django.setup()

from apps.presentation_studio.handlers.markdown_slide_parser import parse_markdown_file
from apps.presentation_studio.handlers.markdown_to_slides_handler import MarkdownToSlidesHandler
from pptx import Presentation

def test_markdown_parsing():
    """Test markdown parsing"""
    print("=" * 60)
    print("Testing Markdown Slide Parser")
    print("=" * 60)
    
    markdown_file = r"C:\Users\achim\github\bfagent\docs\3_Neurowissenschaft_Slides_Detailliert_Teil1.md"
    
    if not os.path.exists(markdown_file):
        print(f"ERROR: File not found: {markdown_file}")
        return False
    
    print(f"\nParsing: {markdown_file}")
    
    try:
        parser = parse_markdown_file(markdown_file)
        print(f"✓ Successfully parsed {len(parser.slides)} slides\n")
        
        # Show first 3 slides
        for slide in parser.slides[:3]:
            print(f"Slide {slide.slide_number}: {slide.title}")
            if slide.headline:
                print(f"  Headline: {slide.headline}")
            print(f"  Content blocks: {len(slide.content_blocks)}")
            if slide.quote:
                print(f"  Quote: {slide.quote[:50]}...")
            print()
        
        return parser
    
    except Exception as e:
        print(f"ERROR: {e}")
        import traceback
        traceback.print_exc()
        return None


def test_slide_creation(parser):
    """Test slide creation"""
    if not parser:
        return False
    
    print("=" * 60)
    print("Testing Slide Creation")
    print("=" * 60)
    
    try:
        # Create presentation
        prs = Presentation()
        print(f"✓ Created presentation with {len(prs.slide_layouts)} layouts")
        
        # Create handler
        handler = MarkdownToSlidesHandler(prs)
        print("✓ Created handler")
        
        # Create first 3 slides
        print("\nCreating slides 1-3...")
        created = handler.create_slides_from_markdown(parser, start_slide=1, end_slide=3)
        print(f"✓ Created {len(created)} slides")
        
        # Save
        output_file = r"C:\Users\achim\github\bfagent\docs\Neurowissenschaft_Test.pptx"
        prs.save(output_file)
        print(f"✓ Saved to: {output_file}")
        
        return True
    
    except Exception as e:
        print(f"ERROR: {e}")
        import traceback
        traceback.print_exc()
        return False


if __name__ == "__main__":
    print("\n" + "=" * 60)
    print("MARKDOWN TO SLIDES TEST")
    print("=" * 60 + "\n")
    
    # Test parsing
    parser = test_markdown_parsing()
    
    if parser:
        # Test slide creation
        success = test_slide_creation(parser)
        
        if success:
            print("\n" + "=" * 60)
            print("✓ ALL TESTS PASSED!")
            print("=" * 60 + "\n")
        else:
            print("\n❌ Slide creation failed")
    else:
        print("\n❌ Parsing failed")
