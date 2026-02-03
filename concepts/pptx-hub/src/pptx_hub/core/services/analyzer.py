"""
Slide Analyzer Service.

Analyzes PowerPoint presentations for structure and content statistics.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import BinaryIO

import structlog
from pptx import Presentation as PptxPresentation

logger = structlog.get_logger(__name__)


@dataclass
class SlideAnalysis:
    """Analysis results for a single slide."""
    
    number: int
    title: str | None
    text_count: int
    word_count: int
    character_count: int
    shape_count: int
    has_images: bool
    has_tables: bool
    has_charts: bool
    has_notes: bool
    layout_name: str | None


@dataclass
class PresentationAnalysis:
    """Analysis results for an entire presentation."""
    
    filename: str
    slide_count: int
    total_text_count: int
    total_word_count: int
    total_character_count: int
    
    # Aggregate statistics
    slides_with_images: int
    slides_with_tables: int
    slides_with_charts: int
    slides_with_notes: int
    
    # Layout breakdown
    layouts_used: dict[str, int]
    
    # Per-slide analysis
    slides: list[SlideAnalysis]
    
    # Metadata
    title: str | None
    author: str | None
    subject: str | None


class SlideAnalyzer:
    """
    Analyze PowerPoint presentations.
    
    This service provides detailed analysis of presentation structure,
    content statistics, and metadata.
    
    Example:
        analyzer = SlideAnalyzer()
        analysis = analyzer.analyze("presentation.pptx")
        
        print(f"Slides: {analysis.slide_count}")
        print(f"Words: {analysis.total_word_count}")
    """
    
    def __init__(self) -> None:
        """Initialize the analyzer."""
        self.log = logger.bind(service="SlideAnalyzer")
    
    def analyze(self, source: str | Path | BinaryIO) -> PresentationAnalysis:
        """
        Analyze a PowerPoint presentation.
        
        Args:
            source: Path to PPTX file or file-like object
            
        Returns:
            PresentationAnalysis with detailed statistics
        """
        self.log.info("analysis_started", source=str(source))
        
        try:
            prs = PptxPresentation(source)
        except Exception as e:
            self.log.error("analysis_failed", error=str(e))
            raise ValueError(f"Failed to open presentation: {e}") from e
        
        # Extract metadata
        filename = str(source) if isinstance(source, (str, Path)) else "uploaded.pptx"
        title = prs.core_properties.title
        author = prs.core_properties.author
        subject = prs.core_properties.subject
        
        # Analyze slides
        slides: list[SlideAnalysis] = []
        layouts_used: dict[str, int] = {}
        
        total_text_count = 0
        total_word_count = 0
        total_character_count = 0
        slides_with_images = 0
        slides_with_tables = 0
        slides_with_charts = 0
        slides_with_notes = 0
        
        for idx, pptx_slide in enumerate(prs.slides, start=1):
            slide_analysis = self._analyze_slide(pptx_slide, idx)
            slides.append(slide_analysis)
            
            # Accumulate totals
            total_text_count += slide_analysis.text_count
            total_word_count += slide_analysis.word_count
            total_character_count += slide_analysis.character_count
            
            if slide_analysis.has_images:
                slides_with_images += 1
            if slide_analysis.has_tables:
                slides_with_tables += 1
            if slide_analysis.has_charts:
                slides_with_charts += 1
            if slide_analysis.has_notes:
                slides_with_notes += 1
            
            # Track layouts
            if slide_analysis.layout_name:
                layouts_used[slide_analysis.layout_name] = (
                    layouts_used.get(slide_analysis.layout_name, 0) + 1
                )
        
        analysis = PresentationAnalysis(
            filename=Path(filename).name,
            slide_count=len(slides),
            total_text_count=total_text_count,
            total_word_count=total_word_count,
            total_character_count=total_character_count,
            slides_with_images=slides_with_images,
            slides_with_tables=slides_with_tables,
            slides_with_charts=slides_with_charts,
            slides_with_notes=slides_with_notes,
            layouts_used=layouts_used,
            slides=slides,
            title=title,
            author=author,
            subject=subject,
        )
        
        self.log.info(
            "analysis_completed",
            slides=analysis.slide_count,
            words=analysis.total_word_count,
        )
        
        return analysis
    
    def _analyze_slide(self, pptx_slide, slide_number: int) -> SlideAnalysis:
        """Analyze a single slide."""
        title: str | None = None
        text_count = 0
        word_count = 0
        character_count = 0
        shape_count = 0
        has_images = False
        has_tables = False
        has_charts = False
        has_notes = False
        layout_name: str | None = None
        
        # Get layout name
        try:
            layout_name = pptx_slide.slide_layout.name
        except Exception:
            pass
        
        # Analyze shapes
        for shape in pptx_slide.shapes:
            shape_count += 1
            
            # Check shape types
            if shape.has_table:
                has_tables = True
            
            if hasattr(shape, "chart"):
                has_charts = True
            
            if hasattr(shape, "image"):
                has_images = True
            
            # Analyze text
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_count += 1
                    word_count += len(text.split())
                    character_count += len(text)
                
                # Check for title
                if shape.is_placeholder:
                    try:
                        from pptx.enum.shapes import PP_PLACEHOLDER
                        if shape.placeholder_format.type in (
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                        ):
                            title = text or None
                    except Exception:
                        pass
        
        # Check notes
        if pptx_slide.has_notes_slide:
            notes_slide = pptx_slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    has_notes = True
        
        return SlideAnalysis(
            number=slide_number,
            title=title,
            text_count=text_count,
            word_count=word_count,
            character_count=character_count,
            shape_count=shape_count,
            has_images=has_images,
            has_tables=has_tables,
            has_charts=has_charts,
            has_notes=has_notes,
            layout_name=layout_name,
        )
    
    def get_summary(self, source: str | Path | BinaryIO) -> dict:
        """
        Get a quick summary of a presentation.
        
        Args:
            source: Path to PPTX file or file-like object
            
        Returns:
            Dictionary with key statistics
        """
        analysis = self.analyze(source)
        
        return {
            "filename": analysis.filename,
            "slides": analysis.slide_count,
            "words": analysis.total_word_count,
            "characters": analysis.total_character_count,
            "images": analysis.slides_with_images,
            "tables": analysis.slides_with_tables,
            "charts": analysis.slides_with_charts,
            "layouts": list(analysis.layouts_used.keys()),
        }
