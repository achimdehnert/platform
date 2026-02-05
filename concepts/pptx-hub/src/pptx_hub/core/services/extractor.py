"""
Text Extractor Service.

Extracts text content from PowerPoint presentations.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import BinaryIO

import structlog
from pptx import Presentation as PptxPresentation
from pptx.util import Emu

from pptx_hub.core.models.presentation import (
    ExtractionResult,
    Presentation,
    Slide,
    SlideText,
    TextType,
)

logger = structlog.get_logger(__name__)


class TextExtractor:
    """
    Extract text content from PowerPoint presentations.
    
    This service reads PPTX files and extracts all text content,
    including titles, body text, notes, and table content.
    
    Example:
        extractor = TextExtractor()
        result = extractor.extract("presentation.pptx")
        
        for slide in result.slides:
            print(f"Slide {slide.number}: {slide.title}")
    """
    
    def __init__(self) -> None:
        """Initialize the extractor."""
        self.log = logger.bind(service="TextExtractor")
    
    def extract(self, source: str | Path | BinaryIO) -> ExtractionResult:
        """
        Extract text from a PowerPoint file.
        
        Args:
            source: Path to PPTX file or file-like object
            
        Returns:
            ExtractionResult with presentation data and statistics
        """
        self.log.info("extraction_started", source=str(source))
        
        warnings: list[str] = []
        errors: list[str] = []
        
        try:
            prs = PptxPresentation(source)
        except Exception as e:
            self.log.error("extraction_failed", error=str(e))
            return ExtractionResult(
                presentation=Presentation(filename=str(source)),
                source_path=str(source),
                errors=[f"Failed to open presentation: {e}"],
            )
        
        # Extract metadata
        filename = str(source) if isinstance(source, (str, Path)) else "uploaded.pptx"
        title = prs.core_properties.title
        author = prs.core_properties.author
        subject = prs.core_properties.subject
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Extract slides
        slides: list[Slide] = []
        total_texts = 0
        total_characters = 0
        
        for idx, pptx_slide in enumerate(prs.slides, start=1):
            try:
                slide = self._extract_slide(pptx_slide, idx)
                slides.append(slide)
                total_texts += slide.text_count
                total_characters += len(slide.all_text)
            except Exception as e:
                warnings.append(f"Slide {idx}: {e}")
                self.log.warning("slide_extraction_warning", slide=idx, error=str(e))
        
        presentation = Presentation(
            filename=Path(filename).name,
            title=title,
            author=author,
            subject=subject,
            slides=slides,
            slide_width=slide_width,
            slide_height=slide_height,
        )
        
        result = ExtractionResult(
            presentation=presentation,
            source_path=str(source),
            total_slides=len(slides),
            total_texts=total_texts,
            total_characters=total_characters,
            warnings=warnings,
            errors=errors,
        )
        
        self.log.info(
            "extraction_completed",
            slides=result.total_slides,
            texts=result.total_texts,
            characters=result.total_characters,
        )
        
        return result
    
    def _extract_slide(self, pptx_slide, slide_number: int) -> Slide:
        """Extract content from a single slide."""
        texts: list[SlideText] = []
        title: str | None = None
        has_images = False
        has_tables = False
        has_charts = False
        shape_count = 0
        
        for shape in pptx_slide.shapes:
            shape_count += 1
            
            # Check shape types
            if shape.has_table:
                has_tables = True
                texts.extend(self._extract_table_text(shape.table, shape.shape_id))
            
            if hasattr(shape, "chart"):
                has_charts = True
            
            if hasattr(shape, "image"):
                has_images = True
            
            # Extract text frame
            if shape.has_text_frame:
                shape_texts = self._extract_text_frame(
                    shape.text_frame,
                    shape.shape_id,
                    shape.is_placeholder,
                )
                
                # Check for title
                if shape.is_placeholder:
                    try:
                        from pptx.enum.shapes import PP_PLACEHOLDER
                        if shape.placeholder_format.type in (
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                        ):
                            if shape_texts and not title:
                                title = shape_texts[0].content
                                shape_texts[0] = SlideText(
                                    **{**shape_texts[0].model_dump(), "text_type": TextType.TITLE}
                                )
                    except Exception:
                        pass
                
                texts.extend(shape_texts)
        
        # Extract notes
        notes: str | None = None
        if pptx_slide.has_notes_slide:
            notes_slide = pptx_slide.notes_slide
            if notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text.strip() or None
        
        # Get layout name
        layout_name: str | None = None
        try:
            layout_name = pptx_slide.slide_layout.name
        except Exception:
            pass
        
        return Slide(
            number=slide_number,
            title=title,
            texts=texts,
            notes=notes,
            layout_name=layout_name,
            has_images=has_images,
            has_tables=has_tables,
            has_charts=has_charts,
            shape_count=shape_count,
        )
    
    def _extract_text_frame(
        self,
        text_frame,
        shape_id: int,
        is_placeholder: bool,
    ) -> list[SlideText]:
        """Extract text from a text frame."""
        texts: list[SlideText] = []
        
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                content = run.text
                if not content or not content.strip():
                    continue
                
                # Get formatting
                font = run.font
                text = SlideText(
                    content=content,
                    text_type=TextType.BODY,
                    shape_id=shape_id,
                    paragraph_index=para_idx,
                    run_index=run_idx,
                    font_name=font.name,
                    font_size=font.size.pt if font.size else None,
                    bold=font.bold or False,
                    italic=font.italic or False,
                )
                texts.append(text)
        
        return texts
    
    def _extract_table_text(self, table, shape_id: int) -> list[SlideText]:
        """Extract text from a table."""
        texts: list[SlideText] = []
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                content = cell.text.strip()
                if content:
                    text = SlideText(
                        content=content,
                        text_type=TextType.TABLE,
                        shape_id=shape_id,
                        paragraph_index=row_idx,
                        run_index=col_idx,
                    )
                    texts.append(text)
        
        return texts
