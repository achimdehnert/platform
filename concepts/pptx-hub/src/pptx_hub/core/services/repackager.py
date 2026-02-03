"""
Repackager Service.

Repackages PowerPoint presentations with modified content.
"""

from __future__ import annotations

import copy
from pathlib import Path
from typing import BinaryIO, Callable

import structlog
from pptx import Presentation as PptxPresentation

logger = structlog.get_logger(__name__)


class Repackager:
    """
    Repackage PowerPoint presentations with modifications.
    
    This service takes a source presentation and creates a new version
    with text replacements or other modifications while preserving
    the original formatting and layout.
    
    Example:
        repackager = Repackager()
        repackager.repackage(
            source="original.pptx",
            output="translated.pptx",
            replacements={"Hello": "Hallo", "World": "Welt"}
        )
    """
    
    def __init__(self) -> None:
        """Initialize the repackager."""
        self.log = logger.bind(service="Repackager")
    
    def repackage(
        self,
        source: str | Path | BinaryIO,
        output: str | Path,
        replacements: dict[str, str] | None = None,
        transform_fn: Callable[[str], str] | None = None,
        skip_notes: bool = False,
    ) -> dict:
        """
        Repackage a presentation with text modifications.
        
        Args:
            source: Path to source PPTX file or file-like object
            output: Path for output PPTX file
            replacements: Dictionary of text replacements {old: new}
            transform_fn: Optional function to transform each text
            skip_notes: If True, don't modify speaker notes
            
        Returns:
            Dictionary with statistics about modifications
        """
        self.log.info("repackage_started", source=str(source), output=str(output))
        
        stats = {
            "slides_processed": 0,
            "texts_replaced": 0,
            "shapes_modified": 0,
            "tables_modified": 0,
            "notes_modified": 0,
        }
        
        try:
            prs = PptxPresentation(source)
        except Exception as e:
            self.log.error("repackage_failed", error=str(e))
            raise ValueError(f"Failed to open presentation: {e}") from e
        
        # Process each slide
        for slide in prs.slides:
            stats["slides_processed"] += 1
            
            # Process shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    modified = self._process_text_frame(
                        shape.text_frame,
                        replacements,
                        transform_fn,
                    )
                    if modified:
                        stats["shapes_modified"] += 1
                        stats["texts_replaced"] += modified
                
                if shape.has_table:
                    modified = self._process_table(
                        shape.table,
                        replacements,
                        transform_fn,
                    )
                    if modified:
                        stats["tables_modified"] += 1
                        stats["texts_replaced"] += modified
            
            # Process notes
            if not skip_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    modified = self._process_text_frame(
                        notes_slide.notes_text_frame,
                        replacements,
                        transform_fn,
                    )
                    if modified:
                        stats["notes_modified"] += 1
                        stats["texts_replaced"] += modified
        
        # Save output
        output_path = Path(output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        self.log.info("repackage_completed", **stats)
        return stats
    
    def _process_text_frame(
        self,
        text_frame,
        replacements: dict[str, str] | None,
        transform_fn: Callable[[str], str] | None,
    ) -> int:
        """Process text in a text frame. Returns count of modifications."""
        modifications = 0
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                original = run.text
                if not original:
                    continue
                
                new_text = original
                
                # Apply replacements
                if replacements:
                    for old, new in replacements.items():
                        if old in new_text:
                            new_text = new_text.replace(old, new)
                
                # Apply transform function
                if transform_fn:
                    new_text = transform_fn(new_text)
                
                # Update if changed
                if new_text != original:
                    run.text = new_text
                    modifications += 1
        
        return modifications
    
    def _process_table(
        self,
        table,
        replacements: dict[str, str] | None,
        transform_fn: Callable[[str], str] | None,
    ) -> int:
        """Process text in a table. Returns count of modifications."""
        modifications = 0
        
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    modified = self._process_text_frame(
                        cell.text_frame,
                        replacements,
                        transform_fn,
                    )
                    modifications += modified
        
        return modifications
    
    def copy_with_replacements(
        self,
        source: str | Path,
        output: str | Path,
        slide_replacements: dict[int, dict[str, str]],
    ) -> dict:
        """
        Copy presentation with per-slide replacements.
        
        Args:
            source: Path to source PPTX file
            output: Path for output PPTX file
            slide_replacements: {slide_number: {old: new}} mapping
            
        Returns:
            Dictionary with statistics
        """
        self.log.info(
            "copy_with_replacements_started",
            source=str(source),
            slides=list(slide_replacements.keys()),
        )
        
        stats = {
            "slides_processed": 0,
            "slides_modified": 0,
            "texts_replaced": 0,
        }
        
        prs = PptxPresentation(source)
        
        for idx, slide in enumerate(prs.slides, start=1):
            stats["slides_processed"] += 1
            
            replacements = slide_replacements.get(idx, {})
            if not replacements:
                continue
            
            slide_modified = False
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    modified = self._process_text_frame(
                        shape.text_frame,
                        replacements,
                        None,
                    )
                    if modified:
                        slide_modified = True
                        stats["texts_replaced"] += modified
            
            if slide_modified:
                stats["slides_modified"] += 1
        
        output_path = Path(output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        self.log.info("copy_with_replacements_completed", **stats)
        return stats
