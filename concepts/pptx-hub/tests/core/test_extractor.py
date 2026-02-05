"""
Tests for text extractor service.
"""

from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from pptx_hub.core.services import TextExtractor


def create_test_pptx() -> BytesIO:
    """Create a test PPTX file in memory."""
    prs = PptxPresentation()
    
    # Add title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Test Presentation"
    slide.placeholders[1].text = "Subtitle text"
    
    # Add content slide
    content_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_layout)
    slide2.shapes.title.text = "Content Slide"
    body = slide2.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


class TestTextExtractor:
    """Tests for TextExtractor."""
    
    def test_extract_from_memory(self):
        """Test extraction from in-memory PPTX."""
        pptx_file = create_test_pptx()
        
        extractor = TextExtractor()
        result = extractor.extract(pptx_file)
        
        assert result.success
        assert result.total_slides == 2
        assert result.total_texts > 0
    
    def test_extract_slide_titles(self):
        """Test that slide titles are extracted."""
        pptx_file = create_test_pptx()
        
        extractor = TextExtractor()
        result = extractor.extract(pptx_file)
        
        assert len(result.slides) == 2
        assert result.slides[0].title == "Test Presentation"
        assert result.slides[1].title == "Content Slide"
    
    def test_extract_body_text(self):
        """Test that body text is extracted."""
        pptx_file = create_test_pptx()
        
        extractor = TextExtractor()
        result = extractor.extract(pptx_file)
        
        # Check content slide has bullet points
        content_slide = result.slides[1]
        all_text = content_slide.all_text
        
        assert "First bullet point" in all_text
        assert "Second bullet point" in all_text
    
    def test_extract_invalid_file(self):
        """Test extraction of invalid file."""
        invalid_file = BytesIO(b"not a pptx file")
        
        extractor = TextExtractor()
        result = extractor.extract(invalid_file)
        
        assert not result.success
        assert len(result.errors) > 0
    
    def test_extraction_result_statistics(self):
        """Test that extraction statistics are calculated."""
        pptx_file = create_test_pptx()
        
        extractor = TextExtractor()
        result = extractor.extract(pptx_file)
        
        assert result.total_slides == 2
        assert result.total_texts > 0
        assert result.total_characters > 0
